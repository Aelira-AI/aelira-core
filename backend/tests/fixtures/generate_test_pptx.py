#!/usr/bin/env python3
"""
Generate synthetic PowerPoint test fixtures for automated testing.

Creates 3 PowerPoint presentations with varying accessibility issues:
1. Lecture deck (moderate contrast issues, mixed alt text)
2. Dark theme (high contrast, but some issues with chart text)
3. Image-heavy presentation (many missing alt text)
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from PIL import Image, ImageDraw, ImageFont
from pathlib import Path
import os


def create_test_image(
    filename: str,
    width: int = 800,
    height: int = 600,
    text: str = "Sample Image",
    bg_color="white",
    text_color="black",
):
    """Create a test image with specified colors."""
    img = Image.new("RGB", (width, height), color=bg_color)
    draw = ImageDraw.Draw(img)

    # Draw border
    draw.rectangle([20, 20, width - 20, height - 20], outline=text_color, width=3)

    # Add text
    try:
        font = ImageFont.truetype(
            "/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf", 48
        )
    except Exception:
        font = ImageFont.load_default()

    # Center text
    bbox = draw.textbbox((0, 0), text, font=font)
    text_width = bbox[2] - bbox[0]
    text_height = bbox[3] - bbox[1]
    x = (width - text_width) // 2
    y = (height - text_height) // 2
    draw.text((x, y), text, fill=text_color, font=font)

    img.save(filename)
    return filename


def generate_lecture_deck():
    """Generate lecture presentation with contrast issues."""
    output_path = Path(__file__).parent / "powerpoint" / "lecture_deck.pptx"
    output_path.parent.mkdir(parents=True, exist_ok=True)

    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Slide 1: Title slide (good contrast)
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout

    # Title
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = "Introduction to Machine Learning"
    p = title_frame.paragraphs[0]
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 0, 0)  # Black text
    p.alignment = PP_ALIGN.CENTER

    # Subtitle
    subtitle_box = slide.shapes.add_textbox(
        Inches(1), Inches(4), Inches(8), Inches(0.75)
    )
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = "CS 450 - Lecture 1"
    p = subtitle_frame.paragraphs[0]
    p.font.size = Pt(28)
    p.font.color.rgb = RGBColor(51, 51, 51)  # Dark gray
    p.alignment = PP_ALIGN.CENTER

    # Background (white)
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(255, 255, 255)

    # Slide 2: Learning objectives (low contrast issue)
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Background (light blue)
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(173, 216, 230)  # Light blue

    # Title (yellow text - LOW CONTRAST)
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.5), Inches(9), Inches(0.75)
    )
    title_frame = title_box.text_frame
    title_frame.text = "Learning Objectives"
    p = title_frame.paragraphs[0]
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 102)  # Light yellow - WCAG FAIL

    # Content (dark text - good contrast)
    content_box = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(5))
    content_frame = content_box.text_frame
    content_frame.text = "• Understand supervised vs unsupervised learning\n"
    content_frame.text += "• Implement linear regression\n"
    content_frame.text += "• Evaluate model performance\n"
    content_frame.text += "• Apply cross-validation techniques"
    for paragraph in content_frame.paragraphs:
        paragraph.font.size = Pt(24)
        paragraph.font.color.rgb = RGBColor(0, 0, 0)
        paragraph.space_before = Pt(12)

    # Slide 3: Key concepts with image (missing alt text)
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Background (white)
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(255, 255, 255)

    # Title
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.5), Inches(9), Inches(0.75)
    )
    title_frame = title_box.text_frame
    title_frame.text = "Supervised Learning Pipeline"
    p = title_frame.paragraphs[0]
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 51, 102)

    # Add image (WITHOUT alt text - accessibility issue)
    img_path = "/tmp/ml_pipeline.png"
    create_test_image(img_path, text="ML Pipeline Diagram")
    slide.shapes.add_picture(img_path, Inches(2), Inches(2), width=Inches(6))
    # NOTE: Not setting alt text - this is the test!

    # Slide 4: Code example (good contrast)
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Background (dark gray)
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(40, 40, 40)

    # Title
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.5), Inches(9), Inches(0.75)
    )
    title_frame = title_box.text_frame
    title_frame.text = "Linear Regression Example"
    p = title_frame.paragraphs[0]
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)

    # Code (monospace, good contrast)
    code_box = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(5))
    code_frame = code_box.text_frame
    code_frame.text = "from sklearn.linear_model import LinearRegression\n\n"
    code_frame.text += "# Create and train model\n"
    code_frame.text += "model = LinearRegression()\n"
    code_frame.text += "model.fit(X_train, y_train)\n\n"
    code_frame.text += "# Make predictions\n"
    code_frame.text += "predictions = model.predict(X_test)"

    for paragraph in code_frame.paragraphs:
        paragraph.font.size = Pt(18)
        paragraph.font.name = "Courier New"
        paragraph.font.color.rgb = RGBColor(220, 220, 220)

    # Slide 5: Summary (medium contrast issue)
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Background (medium gray)
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(128, 128, 128)  # Medium gray

    # Title (light gray text - LOW CONTRAST)
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.5), Inches(9), Inches(0.75)
    )
    title_frame = title_box.text_frame
    title_frame.text = "Key Takeaways"
    p = title_frame.paragraphs[0]
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = RGBColor(192, 192, 192)  # Light gray - WCAG FAIL

    # Content
    content_box = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(5))
    content_frame = content_box.text_frame
    content_frame.text = "1. ML requires clean, labeled data\n"
    content_frame.text += "2. Always split data into train/test sets\n"
    content_frame.text += "3. Evaluate using appropriate metrics\n"
    content_frame.text += "4. Iterate and improve"
    for paragraph in content_frame.paragraphs:
        paragraph.font.size = Pt(24)
        paragraph.font.color.rgb = RGBColor(255, 255, 255)
        paragraph.space_before = Pt(12)

    prs.save(str(output_path))
    print(f"✅ Generated: {output_path}")

    # Cleanup
    if os.path.exists(img_path):
        os.remove(img_path)


def generate_dark_theme():
    """Generate dark theme presentation with subtle contrast issues."""
    output_path = Path(__file__).parent / "powerpoint" / "dark_theme.pptx"
    output_path.parent.mkdir(parents=True, exist_ok=True)

    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Slide 1: Title (high contrast)
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Background (black)
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0, 0, 0)

    # Title
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = "Cybersecurity Fundamentals"
    p = title_frame.paragraphs[0]
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 255, 127)  # Spring green
    p.alignment = PP_ALIGN.CENTER

    # Subtitle
    subtitle_box = slide.shapes.add_textbox(
        Inches(1), Inches(4), Inches(8), Inches(0.75)
    )
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = "Module 3: Network Security"
    p = subtitle_frame.paragraphs[0]
    p.font.size = Pt(28)
    p.font.color.rgb = RGBColor(200, 200, 200)
    p.alignment = PP_ALIGN.CENTER

    # Slide 2: Content with chart (chart text too small/low contrast)
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Background (dark blue)
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0, 32, 64)

    # Title
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.5), Inches(9), Inches(0.75)
    )
    title_frame = title_box.text_frame
    title_frame.text = "Common Attack Vectors"
    p = title_frame.paragraphs[0]
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)

    # Add chart image with good alt text (for comparison)
    img_path = "/tmp/attack_vectors.png"
    create_test_image(
        img_path, text="Attack Vectors", bg_color="lightblue", text_color="darkblue"
    )
    pic = slide.shapes.add_picture(img_path, Inches(1), Inches(2), width=Inches(5))
    # THIS ONE HAS ALT TEXT (for comparison)
    pic._element.nvPicPr.cNvPr.set(
        "descr",
        "Bar chart showing phishing (45%), malware (30%), SQL injection (15%), DDoS (10%)",
    )

    # Content
    content_box = slide.shapes.add_textbox(Inches(6.5), Inches(2), Inches(3), Inches(4))
    content_frame = content_box.text_frame
    content_frame.text = "Key Statistics:\n\n"
    content_frame.text += "• Phishing: 45%\n"
    content_frame.text += "• Malware: 30%\n"
    content_frame.text += "• SQL Injection: 15%\n"
    content_frame.text += "• DDoS: 10%"
    for paragraph in content_frame.paragraphs:
        paragraph.font.size = Pt(20)
        paragraph.font.color.rgb = RGBColor(255, 255, 255)

    # Slide 3: Two-column layout (one with contrast issue)
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Background (dark gray)
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(32, 32, 32)

    # Title
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.5), Inches(9), Inches(0.75)
    )
    title_frame = title_box.text_frame
    title_frame.text = "Defense Strategies"
    p = title_frame.paragraphs[0]
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)

    # Left column (good contrast)
    left_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(1.5), Inches(4.5), Inches(5)
    )
    left_frame = left_box.text_frame
    left_frame.text = "Prevention:\n\n"
    left_frame.text += "• Strong authentication\n"
    left_frame.text += "• Regular updates\n"
    left_frame.text += "• Employee training\n"
    left_frame.text += "• Network segmentation"
    for paragraph in left_frame.paragraphs:
        paragraph.font.size = Pt(20)
        paragraph.font.color.rgb = RGBColor(0, 255, 0)  # Bright green

    # Right column (medium contrast - BORDERLINE)
    right_box = slide.shapes.add_textbox(Inches(5.5), Inches(1.5), Inches(4), Inches(5))
    right_frame = right_box.text_frame
    right_frame.text = "Detection:\n\n"
    right_frame.text += "• SIEM systems\n"
    right_frame.text += "• IDS/IPS\n"
    right_frame.text += "• Log analysis\n"
    right_frame.text += "• Anomaly detection"
    for paragraph in right_frame.paragraphs:
        paragraph.font.size = Pt(20)
        paragraph.font.color.rgb = RGBColor(100, 100, 100)  # Gray - BORDERLINE CONTRAST

    # Slide 4: Code/terminal example (good contrast)
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Background (black)
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0, 0, 0)

    # Title
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.5), Inches(9), Inches(0.75)
    )
    title_frame = title_box.text_frame
    title_frame.text = "Port Scanning Example"
    p = title_frame.paragraphs[0]
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 255, 0)

    # Terminal output
    terminal_box = slide.shapes.add_textbox(
        Inches(1), Inches(1.5), Inches(8), Inches(5)
    )
    terminal_frame = terminal_box.text_frame
    terminal_frame.text = "$ nmap -sS -p 1-1000 192.168.1.1\n\n"
    terminal_frame.text += "Starting Nmap 7.94\n"
    terminal_frame.text += "Nmap scan report for router.local (192.168.1.1)\n"
    terminal_frame.text += "PORT    STATE SERVICE\n"
    terminal_frame.text += "22/tcp  open  ssh\n"
    terminal_frame.text += "80/tcp  open  http\n"
    terminal_frame.text += "443/tcp open  https"

    for paragraph in terminal_frame.paragraphs:
        paragraph.font.size = Pt(18)
        paragraph.font.name = "Courier New"
        paragraph.font.color.rgb = RGBColor(0, 255, 0)

    prs.save(str(output_path))
    print(f"✅ Generated: {output_path}")

    # Cleanup
    if os.path.exists(img_path):
        os.remove(img_path)


def generate_image_heavy():
    """Generate image-heavy presentation with many missing alt texts."""
    output_path = Path(__file__).parent / "powerpoint" / "image_heavy.pptx"
    output_path.parent.mkdir(parents=True, exist_ok=True)

    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Slide 1: Title
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(255, 255, 255)

    title_box = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(8), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = "Photography Portfolio 2025"
    p = title_frame.paragraphs[0]
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 0, 0)
    p.alignment = PP_ALIGN.CENTER

    # Slide 2: 4 images, only 1 has alt text
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(245, 245, 245)

    # Title
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.3), Inches(9), Inches(0.5)
    )
    title_frame = title_box.text_frame
    title_frame.text = "Nature Collection"
    p = title_frame.paragraphs[0]
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 0, 0)

    # Image 1 (top-left) - NO ALT TEXT
    img1_path = "/tmp/nature1.png"
    create_test_image(img1_path, 400, 300, "Landscape", "lightgreen", "darkgreen")
    slide.shapes.add_picture(img1_path, Inches(0.5), Inches(1.2), width=Inches(4.5))

    # Image 2 (top-right) - HAS ALT TEXT
    img2_path = "/tmp/nature2.png"
    create_test_image(img2_path, 400, 300, "Sunset", "lightyellow", "darkorange")
    pic2 = slide.shapes.add_picture(
        img2_path, Inches(5.25), Inches(1.2), width=Inches(4.5)
    )
    pic2._element.nvPicPr.cNvPr.set(
        "descr", "Vibrant sunset over mountain range with orange and purple clouds"
    )

    # Image 3 (bottom-left) - NO ALT TEXT
    img3_path = "/tmp/nature3.png"
    create_test_image(img3_path, 400, 300, "Forest", "lightblue", "darkblue")
    slide.shapes.add_picture(img3_path, Inches(0.5), Inches(4.2), width=Inches(4.5))

    # Image 4 (bottom-right) - NO ALT TEXT
    img4_path = "/tmp/nature4.png"
    create_test_image(img4_path, 400, 300, "Ocean", "lightcoral", "darkred")
    slide.shapes.add_picture(img4_path, Inches(5.25), Inches(4.2), width=Inches(4.5))

    # Slide 3: Single large image - NO ALT TEXT
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0, 0, 0)

    img5_path = "/tmp/feature.png"
    create_test_image(img5_path, 800, 600, "Featured Work", "gray", "white")
    slide.shapes.add_picture(img5_path, Inches(1), Inches(0.75), width=Inches(8))

    # Slide 4: Image with caption - NO ALT TEXT
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(255, 255, 255)

    # Title
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.5), Inches(9), Inches(0.75)
    )
    title_frame = title_box.text_frame
    title_frame.text = "Urban Architecture"
    p = title_frame.paragraphs[0]
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 0, 0)

    # Image
    img6_path = "/tmp/architecture.png"
    create_test_image(img6_path, 600, 400, "Skyscraper", "silver", "black")
    slide.shapes.add_picture(img6_path, Inches(1.5), Inches(1.5), width=Inches(7))

    # Caption (but no alt text on image!)
    caption_box = slide.shapes.add_textbox(
        Inches(1.5), Inches(6), Inches(7), Inches(0.75)
    )
    caption_frame = caption_box.text_frame
    caption_frame.text = "Downtown skyline at golden hour"
    p = caption_frame.paragraphs[0]
    p.font.size = Pt(18)
    p.font.italic = True
    p.font.color.rgb = RGBColor(100, 100, 100)
    p.alignment = PP_ALIGN.CENTER

    # Slide 5: Chart without alt text
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(255, 255, 255)

    # Title
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.5), Inches(9), Inches(0.75)
    )
    title_frame = title_box.text_frame
    title_frame.text = "Exhibition Statistics"
    p = title_frame.paragraphs[0]
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 0, 0)

    # Chart image - NO ALT TEXT
    img7_path = "/tmp/stats_chart.png"
    create_test_image(img7_path, 700, 500, "Visitor Stats", "white", "black")
    slide.shapes.add_picture(img7_path, Inches(1.5), Inches(1.5), width=Inches(7))

    prs.save(str(output_path))
    print(f"✅ Generated: {output_path}")

    # Cleanup
    for img_path in [
        img1_path,
        img2_path,
        img3_path,
        img4_path,
        img5_path,
        img6_path,
        img7_path,
    ]:
        if os.path.exists(img_path):
            os.remove(img_path)


def main():
    """Generate all PowerPoint test fixtures."""
    print("🔄 Generating PowerPoint test fixtures...\n")

    generate_lecture_deck()
    generate_dark_theme()
    generate_image_heavy()

    print("\n✅ All PowerPoint fixtures generated successfully!")
    print(f"📁 Location: {Path(__file__).parent / 'powerpoint'}")


if __name__ == "__main__":
    main()
