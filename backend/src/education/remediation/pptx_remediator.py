"""
PowerPoint Remediator for Aelira Auto-Remediation Engine.

This module provides automatic remediation for accessibility issues in
Microsoft PowerPoint presentations (.pptx files).

Supported auto-fixes:
- Add/update alt text for images and shapes
- Fix color contrast issues
- Add slide titles
- Fix reading order
- Add speaker notes for complex visuals
"""

import logging
from typing import Any, Dict, List, Optional, Tuple
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn

from .base import (
    BaseRemediator,
    RemediationIssue,
    IssueCategory,
    IssueSeverity,
    RemediationConfig,
)

logger = logging.getLogger(__name__)


class PptxRemediator(BaseRemediator):
    """
    Remediator for Microsoft PowerPoint presentations (.pptx).

    Automatically fixes accessibility issues including:
    - Missing or inadequate alt text on images and shapes
    - Insufficient color contrast
    - Missing slide titles
    - Incorrect reading order
    - Complex visuals without text descriptions

    Usage:
        issues = [{'type': 'alt_text', 'severity': 'high', ...}]
        remediator = PptxRemediator('presentation.pptx', issues)
        result = remediator.remediate()
    """

    DOCUMENT_TYPE = "powerpoint"
    SUPPORTED_EXTENSIONS = [".pptx"]

    AUTO_FIXABLE_CATEGORIES = [
        IssueCategory.ALT_TEXT,
        IssueCategory.CONTRAST,
        IssueCategory.STRUCTURE,
        IssueCategory.READING_ORDER,
    ]

    # Minimum contrast ratio for WCAG AA
    MIN_CONTRAST_NORMAL = 4.5
    MIN_CONTRAST_LARGE = 3.0

    def __init__(
        self,
        file_path: str,
        issues: List[Dict[str, Any]],
        config: Optional[RemediationConfig] = None,
        ai_client: Optional[Any] = None,
    ):
        """Initialize the PowerPoint remediator."""
        super().__init__(file_path, issues, config, ai_client)
        self._presentation: Optional[Presentation] = None

    def _load_document(self) -> Presentation:
        """Load the PowerPoint presentation for editing."""
        logger.info(f"Loading PowerPoint: {self.file_path}")
        self._presentation = Presentation(self.file_path)
        return self._presentation

    def _save_document(self, document: Presentation) -> str:
        """Save the remediated PowerPoint presentation."""
        output_path = self._get_output_path()
        logger.info(f"Saving remediated presentation to: {output_path}")

        # Ensure output directory exists
        Path(output_path).parent.mkdir(parents=True, exist_ok=True)

        document.save(output_path)
        return output_path

    def can_auto_fix(self, issue: RemediationIssue) -> bool:
        """
        Determine if an issue can be automatically fixed.

        Auto-fixable issues:
        - Alt text: If we have AI or can generate placeholder
        - Contrast: If we can calculate and adjust colors
        - Structure: If slide title is missing
        - Reading order: If we can reorder shapes
        """
        if issue.category not in self.AUTO_FIXABLE_CATEGORIES:
            return False

        # Category-specific checks
        if issue.category == IssueCategory.ALT_TEXT:
            # Can fix with AI, or with placeholder
            return self.config.use_ai or self.config.fix_alt_text

        if issue.category == IssueCategory.CONTRAST:
            # Can fix if we have color information
            return bool(
                issue.metadata.get("foreground_color")
                or issue.metadata.get("background_color")
            )

        if issue.category == IssueCategory.STRUCTURE:
            # Can add slide title
            return issue.metadata.get("issue_type") == "missing_title"

        if issue.category == IssueCategory.READING_ORDER:
            # Can adjust reading order
            return bool(issue.metadata.get("slide_index") is not None)

        return False

    def apply_fix(
        self, issue: RemediationIssue, document: Presentation, fix_content: str
    ) -> bool:
        """
        Apply a fix to the PowerPoint presentation.

        Args:
            issue: The issue being fixed
            document: The PowerPoint presentation object
            fix_content: The content to apply as the fix

        Returns:
            True if fix was applied successfully
        """
        try:
            if issue.category == IssueCategory.ALT_TEXT:
                return self._apply_alt_text_fix(issue, document, fix_content)

            if issue.category == IssueCategory.CONTRAST:
                return self._apply_contrast_fix(issue, document, fix_content)

            if issue.category == IssueCategory.STRUCTURE:
                return self._apply_structure_fix(issue, document, fix_content)

            if issue.category == IssueCategory.READING_ORDER:
                return self._apply_reading_order_fix(issue, document, fix_content)

            return False

        except Exception as e:
            logger.error(f"Failed to apply fix for issue {issue.id}: {e}")
            return False

    def _apply_alt_text_fix(
        self, issue: RemediationIssue, document: Presentation, alt_text: str
    ) -> bool:
        """Apply alt text fix to a shape or image."""
        try:
            slide_index = issue.metadata.get("slide_index")
            shape_index = issue.metadata.get("shape_index")
            shape_name = issue.metadata.get("shape_name", "")

            if slide_index is None:
                logger.warning(f"No slide index for alt text fix: {issue.id}")
                return False

            if slide_index >= len(document.slides):
                logger.warning(f"Slide index out of range: {slide_index}")
                return False

            slide = document.slides[slide_index]

            # Find the shape
            target_shape = None

            if shape_index is not None and shape_index < len(slide.shapes):
                target_shape = slide.shapes[shape_index]
            elif shape_name:
                # Find by name
                for shape in slide.shapes:
                    if shape.name == shape_name:
                        target_shape = shape
                        break

            if target_shape is None:
                # Try to find any image without alt text
                for shape in slide.shapes:
                    if self._is_image_shape(shape):
                        existing_alt = self._get_shape_alt_text(shape)
                        if not existing_alt or existing_alt.strip() == "":
                            target_shape = shape
                            break

            if target_shape is None:
                logger.warning(
                    f"Could not find shape for alt text fix on slide {slide_index}"
                )
                return False

            # Set alt text
            self._set_shape_alt_text(target_shape, alt_text)
            logger.info(f"Applied alt text to shape on slide {slide_index + 1}")
            return True

        except Exception as e:
            logger.error(f"Error applying alt text fix: {e}")
            return False

    def _is_image_shape(self, shape) -> bool:
        """Check if a shape is an image or picture."""
        try:
            if hasattr(shape, "shape_type"):
                return shape.shape_type == MSO_SHAPE_TYPE.PICTURE
            return False
        except Exception:
            return False

    def _get_shape_alt_text(self, shape) -> str:
        """Get the alt text of a shape."""
        try:
            # Access the alternative text through the shape's XML
            nvSpPr = shape._element.find(".//" + qn("p:nvSpPr"))
            if nvSpPr is not None:
                cNvPr = nvSpPr.find(qn("p:cNvPr"))
                if cNvPr is not None:
                    return cNvPr.get("descr", "")

            # Try through nvPicPr for pictures
            nvPicPr = shape._element.find(".//" + qn("p:nvPicPr"))
            if nvPicPr is not None:
                cNvPr = nvPicPr.find(qn("p:cNvPr"))
                if cNvPr is not None:
                    return cNvPr.get("descr", "")

            return ""
        except Exception:
            return ""

    def _set_shape_alt_text(self, shape, alt_text: str):
        """Set the alt text of a shape."""
        try:
            # Find cNvPr element (common non-visual properties)
            # Try different paths for different shape types

            cNvPr = None

            # For pictures
            nvPicPr = shape._element.find(".//" + qn("p:nvPicPr"))
            if nvPicPr is not None:
                cNvPr = nvPicPr.find(qn("p:cNvPr"))

            # For shapes
            if cNvPr is None:
                nvSpPr = shape._element.find(".//" + qn("p:nvSpPr"))
                if nvSpPr is not None:
                    cNvPr = nvSpPr.find(qn("p:cNvPr"))

            # For group shapes
            if cNvPr is None:
                nvGrpSpPr = shape._element.find(".//" + qn("p:nvGrpSpPr"))
                if nvGrpSpPr is not None:
                    cNvPr = nvGrpSpPr.find(qn("p:cNvPr"))

            if cNvPr is not None:
                cNvPr.set("descr", alt_text)
                logger.debug(f"Set alt text on shape: {shape.name}")
            else:
                logger.warning(f"Could not find cNvPr element for shape: {shape.name}")

        except Exception as e:
            logger.error(f"Error setting alt text: {e}")

    def _apply_contrast_fix(
        self, issue: RemediationIssue, document: Presentation, fix_content: str
    ) -> bool:
        """Apply contrast fix by adjusting colors."""
        try:
            slide_index = issue.metadata.get("slide_index")
            shape_index = issue.metadata.get("shape_index")
            fg_color = issue.metadata.get("foreground_color")
            bg_color = issue.metadata.get("background_color")

            if slide_index is None or slide_index >= len(document.slides):
                return False

            slide = document.slides[slide_index]

            # Find the shape with the contrast issue
            if shape_index is not None and shape_index < len(slide.shapes):
                shape = slide.shapes[shape_index]
            else:
                logger.warning("Could not find shape for contrast fix")
                return False

            # Calculate the fix - darken or lighten colors for better contrast
            if fg_color and bg_color:
                new_fg = self._adjust_for_contrast(fg_color, bg_color)
                if new_fg and self._apply_text_color(shape, new_fg):
                    logger.info(f"Applied contrast fix on slide {slide_index + 1}")
                    return True

            return False

        except Exception as e:
            logger.error(f"Error applying contrast fix: {e}")
            return False

    def _adjust_for_contrast(self, fg_hex: str, bg_hex: str) -> Optional[str]:
        """Adjust foreground color for better contrast against background."""
        try:
            # Parse colors
            fg_rgb = self._hex_to_rgb(fg_hex)
            bg_rgb = self._hex_to_rgb(bg_hex)

            if not fg_rgb or not bg_rgb:
                return None

            # Calculate current contrast
            current_ratio = self._calculate_contrast_ratio(fg_rgb, bg_rgb)

            if current_ratio >= self.MIN_CONTRAST_NORMAL:
                return fg_hex  # Already sufficient

            # Determine if we should lighten or darken
            self._relative_luminance(fg_rgb)
            bg_lum = self._relative_luminance(bg_rgb)

            if bg_lum > 0.5:
                # Dark text on light background - make text darker
                new_rgb = self._darken_color(fg_rgb)
            else:
                # Light text on dark background - make text lighter
                new_rgb = self._lighten_color(fg_rgb)

            return self._rgb_to_hex(new_rgb)

        except Exception as e:
            logger.error(f"Error adjusting contrast: {e}")
            return None

    def _hex_to_rgb(self, hex_color: str) -> Optional[Tuple[int, int, int]]:
        """Convert hex color to RGB tuple."""
        try:
            hex_color = hex_color.lstrip("#")
            if len(hex_color) == 6:
                return tuple(int(hex_color[i : i + 2], 16) for i in (0, 2, 4))
            return None
        except Exception:
            return None

    def _rgb_to_hex(self, rgb: Tuple[int, int, int]) -> str:
        """Convert RGB tuple to hex color."""
        return "#{:02x}{:02x}{:02x}".format(*rgb)

    def _relative_luminance(self, rgb: Tuple[int, int, int]) -> float:
        """Calculate relative luminance of a color."""
        r, g, b = [x / 255.0 for x in rgb]

        def adjust(c):
            if c <= 0.03928:
                return c / 12.92
            return ((c + 0.055) / 1.055) ** 2.4

        return 0.2126 * adjust(r) + 0.7152 * adjust(g) + 0.0722 * adjust(b)

    def _calculate_contrast_ratio(
        self, fg: Tuple[int, int, int], bg: Tuple[int, int, int]
    ) -> float:
        """Calculate contrast ratio between two colors."""
        l1 = self._relative_luminance(fg)
        l2 = self._relative_luminance(bg)

        lighter = max(l1, l2)
        darker = min(l1, l2)

        return (lighter + 0.05) / (darker + 0.05)

    def _darken_color(
        self, rgb: Tuple[int, int, int], factor: float = 0.3
    ) -> Tuple[int, int, int]:
        """Darken a color by a factor."""
        return tuple(max(0, int(c * (1 - factor))) for c in rgb)

    def _lighten_color(
        self, rgb: Tuple[int, int, int], factor: float = 0.3
    ) -> Tuple[int, int, int]:
        """Lighten a color by a factor."""
        return tuple(min(255, int(c + (255 - c) * factor)) for c in rgb)

    def _apply_text_color(self, shape, hex_color: str) -> bool:
        """Apply a text color to a shape."""
        try:
            if not hasattr(shape, "text_frame"):
                return False

            rgb = self._hex_to_rgb(hex_color)
            if not rgb:
                return False

            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.color.rgb = RGBColor(*rgb)

            return True
        except Exception as e:
            logger.error(f"Error applying text color: {e}")
            return False

    def _apply_structure_fix(
        self, issue: RemediationIssue, document: Presentation, fix_content: str
    ) -> bool:
        """Apply structure fix (e.g., add missing slide title)."""
        try:
            slide_index = issue.metadata.get("slide_index")
            issue_type = issue.metadata.get("issue_type")

            if slide_index is None or slide_index >= len(document.slides):
                return False

            slide = document.slides[slide_index]

            if issue_type == "missing_title":
                return self._add_slide_title(slide, fix_content, slide_index)

            return False

        except Exception as e:
            logger.error(f"Error applying structure fix: {e}")
            return False

    def _add_slide_title(self, slide, title_text: str, slide_index: int) -> bool:
        """Add a title to a slide."""
        try:
            # Check if slide has a title placeholder
            title_shape = None

            for shape in slide.shapes:
                if shape.is_placeholder:
                    if shape.placeholder_format.type == 1:  # Title placeholder
                        title_shape = shape
                        break

            if title_shape:
                # Use existing title placeholder
                title_shape.text = title_text
            else:
                # Add a text box as title
                left = Inches(0.5)
                top = Inches(0.3)
                width = Inches(9)
                height = Inches(0.8)

                text_box = slide.shapes.add_textbox(left, top, width, height)
                text_frame = text_box.text_frame
                text_frame.text = title_text

                # Format as title
                for paragraph in text_frame.paragraphs:
                    paragraph.font.size = Pt(32)
                    paragraph.font.bold = True

            logger.info(f"Added title to slide {slide_index + 1}: {title_text[:50]}...")
            return True

        except Exception as e:
            logger.error(f"Error adding slide title: {e}")
            return False

    def _apply_reading_order_fix(
        self, issue: RemediationIssue, document: Presentation, fix_content: str
    ) -> bool:
        """Apply reading order fix by reordering shapes."""
        try:
            slide_index = issue.metadata.get("slide_index")
            suggested_order = issue.metadata.get("suggested_order", [])

            if slide_index is None or slide_index >= len(document.slides):
                return False

            if not suggested_order:
                # Try to determine logical order
                suggested_order = self._determine_reading_order(
                    document.slides[slide_index]
                )

            if not suggested_order:
                return False

            slide = document.slides[slide_index]

            # Reorder shapes by manipulating XML
            # This is complex in python-pptx, log for now
            logger.info(f"Reading order fix requested for slide {slide_index + 1}")
            logger.info(f"Suggested order: {suggested_order[:5]}...")

            # Add to speaker notes as guidance (not a structural fix)
            self._add_reading_order_note(slide, suggested_order)
            return False  # Speaker notes don't fix the actual reading order

        except Exception as e:
            logger.error(f"Error applying reading order fix: {e}")
            return False

    def _determine_reading_order(self, slide) -> List[str]:
        """Determine logical reading order for shapes on a slide."""
        shapes_with_pos = []

        for shape in slide.shapes:
            if hasattr(shape, "left") and hasattr(shape, "top"):
                shapes_with_pos.append(
                    {
                        "name": shape.name,
                        "left": shape.left,
                        "top": shape.top,
                    }
                )

        # Sort by position: top-to-bottom, left-to-right
        shapes_with_pos.sort(key=lambda s: (s["top"], s["left"]))

        return [s["name"] for s in shapes_with_pos]

    def _add_reading_order_note(self, slide, order: List[str]):
        """Add reading order information to speaker notes."""
        try:
            notes_slide = slide.notes_slide
            notes_frame = notes_slide.notes_text_frame

            order_text = "\n\n[Accessibility Note - Reading Order]\n"
            order_text += "Suggested reading order:\n"
            for i, name in enumerate(order[:10], 1):
                order_text += f"{i}. {name}\n"

            notes_frame.text += order_text

        except Exception as e:
            logger.warning(f"Could not add reading order note: {e}")

    def _get_rule_based_fix(
        self, issue: RemediationIssue, document: Any
    ) -> Optional[str]:
        """Get a rule-based fix for an issue."""
        if issue.category == IssueCategory.ALT_TEXT:
            # Use pre-generated alt text from the scanner if available
            # Scanner stores as "suggested_alt_text", check both keys
            generated_alt = issue.metadata.get(
                "suggested_alt_text"
            ) or issue.metadata.get("generated_alt_text")
            if generated_alt:
                return generated_alt
            # Use fix_suggestion if available
            if issue.fix_suggestion:
                return issue.fix_suggestion
            # Return None to let AI generation handle it in _generate_fix()
            return None

        if issue.category == IssueCategory.STRUCTURE:
            if issue.metadata.get("issue_type") == "missing_title":
                slide_index = issue.metadata.get("slide_index", 0)
                return f"Slide {slide_index + 1}"

        if issue.category == IssueCategory.CONTRAST:
            # Return adjusted color based on metadata
            fg_color = issue.metadata.get("foreground_color")
            bg_color = issue.metadata.get("background_color")
            if fg_color and bg_color:
                return self._adjust_for_contrast(fg_color, bg_color)

        return None

    def _get_ai_generated_fix(
        self, issue: RemediationIssue, document: Any
    ) -> Optional[str]:
        """Get an AI-generated fix for an issue."""
        if not self.ai_client:
            return None

        try:
            self.result.ai_calls_made += 1

            if issue.category == IssueCategory.ALT_TEXT:
                return self._generate_alt_text_with_ai(issue, document)

            if issue.category == IssueCategory.STRUCTURE:
                return self._generate_title_with_ai(issue, document)

            return None

        except Exception as e:
            logger.error(f"AI generation failed: {e}")
            return None

    def _generate_alt_text_with_ai(
        self, issue: RemediationIssue, document: Presentation
    ) -> Optional[str]:
        """Generate alt text for a shape using AI vision.

        Extracts the actual image from the PPTX shape and sends it to a
        vision model (e.g. Gemini) so the alt text describes the real
        visual content rather than being guessed from surrounding text.
        Falls back to text-based generation if image extraction fails.
        """
        slide_index = issue.metadata.get("slide_index", 0)

        # Gather slide context for the prompt
        context_text = ""
        target_shape = None
        if document and slide_index < len(document.slides):
            slide = document.slides[slide_index]
            texts = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    texts.append(shape.text)
            context_text = " ".join(texts)

            # Locate the target shape to extract image bytes
            shape_name = issue.metadata.get("shape_name", "")
            shape_id = issue.metadata.get("shape_id")
            for shape in slide.shapes:
                if shape_name and shape.name == shape_name:
                    target_shape = shape
                    break
                if shape_id is not None and shape.shape_id == shape_id:
                    target_shape = shape
                    break

        # Try vision-based generation with actual image data
        image_bytes = None
        if target_shape is not None:
            try:
                image_bytes = target_shape.image.blob
            except Exception:
                pass  # Shape may not have an image blob

        if image_bytes and hasattr(self.ai_client, "analyze_image_sync"):
            from ...utils.security import sanitize_for_prompt

            safe_context = (
                sanitize_for_prompt(context_text, max_length=500)
                if context_text
                else "None"
            )
            prompt = (
                "Generate concise, descriptive alt text for this image "
                "from a PowerPoint slide in a higher education context.\n\n"
                f"Slide context: {safe_context}\n\n"
                "Requirements:\n"
                "- Be concise (under 125 characters)\n"
                "- Describe the visual content and its purpose\n"
                "- Don't start with 'Image of' or 'Picture of'\n"
                "- Focus on what matters for understanding the slide\n\n"
                "Generate only the alt text, nothing else:"
            )
            try:
                result = self.ai_client.analyze_image_sync(
                    image_data=image_bytes,
                    prompt=prompt,
                    max_tokens=200,
                )
                if result.get("success") and result.get("content"):
                    alt_text = result["content"].strip().strip("\"'")
                    if alt_text:
                        return alt_text[:125]
            except Exception as e:
                logger.warning(f"Vision-based alt text generation failed: {e}")

        # Fallback: text-based generation using slide context
        if hasattr(self.ai_client, "generate_text_sync"):
            from ...utils.security import sanitize_for_prompt

            safe_context = (
                sanitize_for_prompt(context_text, max_length=500)
                if context_text
                else "None"
            )
            safe_shape = sanitize_for_prompt(
                issue.metadata.get("shape_name", "Unknown"), max_length=100
            )
            prompt = (
                "Generate concise, descriptive alt text for an image "
                "in a PowerPoint presentation.\n\n"
                f"Slide context: {safe_context}\n"
                f"Shape name: {safe_shape}\n\n"
                "Requirements:\n"
                "- Be concise (under 125 characters)\n"
                "- Describe the likely visual content based on context\n"
                "- Don't start with 'Image of' or 'Picture of'\n\n"
                "Generate only the alt text, nothing else:"
            )
            try:
                result = self.ai_client.generate_text_sync(
                    prompt=prompt,
                    max_tokens=200,
                    temperature=0.3,
                )
                if result.get("success") and result.get("content"):
                    alt_text = result["content"].strip().strip("\"'")
                    if alt_text:
                        return alt_text[:125]
            except Exception as e:
                logger.warning(f"Text-based alt text generation failed: {e}")

        return f"Visual element on slide {slide_index + 1}"

    def _generate_title_with_ai(
        self, issue: RemediationIssue, document: Presentation
    ) -> Optional[str]:
        """Generate a slide title using AI."""
        slide_index = issue.metadata.get("slide_index", 0)

        # Get slide content
        context_text = ""
        if document and slide_index < len(document.slides):
            slide = document.slides[slide_index]
            texts = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    texts.append(shape.text)
            context_text = " ".join(texts)

        prompt = (
            "Generate a concise slide title based on the slide content.\n\n"
            f"Slide content:\n{context_text[:500] if context_text else 'No content available'}\n\n"
            "Requirements:\n"
            "- Keep it short (3-8 words)\n"
            "- Make it descriptive of the slide's main topic\n"
            "- Use title case\n\n"
            "Generate only the title, nothing else:"
        )

        try:
            if hasattr(self.ai_client, "generate_text_sync"):
                result = self.ai_client.generate_text_sync(
                    prompt=prompt,
                    max_tokens=100,
                    temperature=0.3,
                )
                if result.get("success") and result.get("content"):
                    return result["content"].strip().strip("\"'")[:100]
        except Exception as e:
            logger.error(f"AI title generation failed: {e}")

        return f"Slide {slide_index + 1}"

    def _calculate_scores(self):
        """Calculate compliance scores for the remediation."""
        if self.result.total_issues > 0:
            severity_penalties = {
                IssueSeverity.CRITICAL: 15,
                IssueSeverity.HIGH: 10,
                IssueSeverity.MEDIUM: 5,
                IssueSeverity.LOW: 2,
            }

            total_penalty = sum(
                severity_penalties.get(issue.severity, 5) for issue in self.issues
            )
            self.result.original_compliance_score = max(0, 100 - total_penalty)

            fixed_penalty_reduction = sum(
                severity_penalties.get(fixed.severity, 5)
                for fixed in self.result.fixed_issues
            )
            remaining_penalty = total_penalty - fixed_penalty_reduction
            self.result.remediated_compliance_score = max(0, 100 - remaining_penalty)

            self.result.improvement = (
                self.result.remediated_compliance_score
                - self.result.original_compliance_score
            )
