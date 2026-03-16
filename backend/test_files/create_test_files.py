#!/usr/bin/env python3
"""
Create test files for cross-scanner integration testing.
Creates PowerPoint with charts and PDF with images/charts.
"""

import os
from pathlib import Path

# Create test directory
TEST_DIR = Path(__file__).parent
TEST_DIR.mkdir(exist_ok=True)


def create_pptx_with_charts():
    """Create a PowerPoint file with various chart types."""
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.chart.data import CategoryChartData
        from pptx.enum.chart import XL_CHART_TYPE

        prs = Presentation()

        # Slide 1: Title slide
        slide_layout = prs.slide_layouts[6]  # Blank layout
        slide = prs.slides.add_slide(slide_layout)

        # Add title text box
        left = Inches(1)
        top = Inches(2)
        width = Inches(8)
        height = Inches(1.5)
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.text = "Q3 2025 Sales Performance Report"
        tf.paragraphs[0].font.size = Pt(44)
        tf.paragraphs[0].font.bold = True

        # Slide 2: Bar chart
        slide = prs.slides.add_slide(slide_layout)
        chart_data = CategoryChartData()
        chart_data.categories = ['North', 'South', 'East', 'West']
        chart_data.add_series('Q1', (19.2, 21.4, 16.7, 28.0))
        chart_data.add_series('Q2', (22.3, 28.6, 15.2, 32.3))
        chart_data.add_series('Q3', (20.4, 26.3, 20.1, 35.9))

        x, y, cx, cy = Inches(1), Inches(1.5), Inches(8), Inches(5)
        chart = slide.shapes.add_chart(
            XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data
        ).chart

        # Add title above chart
        txBox = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(0.8))
        tf = txBox.text_frame
        tf.text = "Regional Sales by Quarter ($ millions)"
        tf.paragraphs[0].font.size = Pt(24)

        # Slide 3: Pie chart
        slide = prs.slides.add_slide(slide_layout)
        chart_data = CategoryChartData()
        chart_data.categories = ['Product A', 'Product B', 'Product C', 'Product D']
        chart_data.add_series('Market Share', (35.2, 28.5, 22.1, 14.2))

        chart = slide.shapes.add_chart(
            XL_CHART_TYPE.PIE, Inches(2), Inches(1.5), Inches(6), Inches(5), chart_data
        ).chart

        txBox = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(0.8))
        tf = txBox.text_frame
        tf.text = "Product Market Share Distribution"
        tf.paragraphs[0].font.size = Pt(24)

        # Slide 4: Line chart
        slide = prs.slides.add_slide(slide_layout)
        chart_data = CategoryChartData()
        chart_data.categories = ['Jan', 'Feb', 'Mar', 'Apr', 'May', 'Jun']
        chart_data.add_series('Revenue', (4.3, 4.8, 5.2, 5.0, 5.9, 6.3))
        chart_data.add_series('Expenses', (3.2, 3.4, 3.8, 3.5, 3.9, 4.1))
        chart_data.add_series('Profit', (1.1, 1.4, 1.4, 1.5, 2.0, 2.2))

        chart = slide.shapes.add_chart(
            XL_CHART_TYPE.LINE, Inches(1), Inches(1.5), Inches(8), Inches(5), chart_data
        ).chart

        txBox = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(0.8))
        tf = txBox.text_frame
        tf.text = "Monthly Financial Trends ($ millions)"
        tf.paragraphs[0].font.size = Pt(24)

        # Save the presentation
        output_path = TEST_DIR / "sales_report_charts.pptx"
        prs.save(str(output_path))
        print(f"✅ Created PowerPoint with charts: {output_path}")
        return output_path

    except ImportError as e:
        print(f"❌ python-pptx not installed: {e}")
        return None


def create_pptx_with_images():
    """Create a PowerPoint file with various images (some with alt text, some without)."""
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt

        prs = Presentation()
        slide_layout = prs.slide_layouts[6]  # Blank layout

        # Slide 1: Image with alt text
        slide = prs.slides.add_slide(slide_layout)
        txBox = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(0.8))
        tf = txBox.text_frame
        tf.text = "Slide with Decorative Background"
        tf.paragraphs[0].font.size = Pt(24)

        # Add decorative image if it exists
        decorative_path = TEST_DIR / "decorative.jpg"
        if decorative_path.exists():
            pic = slide.shapes.add_picture(
                str(decorative_path), Inches(0.5), Inches(1.5), width=Inches(9)
            )
            # Note: python-pptx doesn't directly support alt text, but we can test detection

        # Slide 2: Image that should have alt text (infographic)
        slide = prs.slides.add_slide(slide_layout)
        txBox = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(0.8))
        tf = txBox.text_frame
        tf.text = "Dashboard Analytics"
        tf.paragraphs[0].font.size = Pt(24)

        infographic_path = TEST_DIR / "infographic.jpg"
        if infographic_path.exists():
            pic = slide.shapes.add_picture(
                str(infographic_path), Inches(1), Inches(1.5), width=Inches(8)
            )

        # Slide 3: Chart image
        slide = prs.slides.add_slide(slide_layout)
        txBox = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(0.8))
        tf = txBox.text_frame
        tf.text = "Performance Metrics"
        tf.paragraphs[0].font.size = Pt(24)

        chart_path = TEST_DIR / "bar_chart.png"
        if chart_path.exists():
            pic = slide.shapes.add_picture(
                str(chart_path), Inches(1.5), Inches(1.5), width=Inches(7)
            )

        # Save the presentation
        output_path = TEST_DIR / "mixed_content.pptx"
        prs.save(str(output_path))
        print(f"✅ Created PowerPoint with images: {output_path}")
        return output_path

    except ImportError as e:
        print(f"❌ python-pptx not installed: {e}")
        return None


def create_pdf_with_images():
    """Create a PDF with embedded images for testing."""
    try:
        from reportlab.lib.pagesizes import letter
        from reportlab.platypus import SimpleDocTemplate, Paragraph, Image, Spacer, Table, TableStyle
        from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
        from reportlab.lib.units import inch
        from reportlab.lib import colors

        output_path = TEST_DIR / "report_with_images.pdf"
        doc = SimpleDocTemplate(str(output_path), pagesize=letter)
        styles = getSampleStyleSheet()
        story = []

        # Title
        title_style = ParagraphStyle(
            'CustomTitle',
            parent=styles['Heading1'],
            fontSize=24,
            spaceAfter=30
        )
        story.append(Paragraph("Quarterly Business Report", title_style))
        story.append(Spacer(1, 20))

        # Introduction
        story.append(Paragraph(
            "This report contains analysis of Q3 2025 performance metrics, "
            "including regional sales data, market share distribution, and financial trends.",
            styles['Normal']
        ))
        story.append(Spacer(1, 20))

        # Add chart image if exists
        chart_path = TEST_DIR / "bar_chart.png"
        if chart_path.exists():
            story.append(Paragraph("Figure 1: Sales Performance Dashboard", styles['Heading2']))
            img = Image(str(chart_path), width=5*inch, height=3*inch)
            story.append(img)
            story.append(Spacer(1, 10))
            story.append(Paragraph(
                "The dashboard above shows key performance indicators for Q3 2025.",
                styles['Normal']
            ))
            story.append(Spacer(1, 20))

        # Add infographic if exists
        infographic_path = TEST_DIR / "infographic.jpg"
        if infographic_path.exists():
            story.append(Paragraph("Figure 2: Analytics Overview", styles['Heading2']))
            img = Image(str(infographic_path), width=5*inch, height=3*inch)
            story.append(img)
            story.append(Spacer(1, 10))
            story.append(Paragraph(
                "Analytics dashboard showing user engagement metrics and trends.",
                styles['Normal']
            ))
            story.append(Spacer(1, 20))

        # Add a data table
        story.append(Paragraph("Table 1: Regional Sales Summary", styles['Heading2']))
        data = [
            ['Region', 'Q1', 'Q2', 'Q3', 'YTD Total'],
            ['North', '$19.2M', '$22.3M', '$20.4M', '$61.9M'],
            ['South', '$21.4M', '$28.6M', '$26.3M', '$76.3M'],
            ['East', '$16.7M', '$15.2M', '$20.1M', '$52.0M'],
            ['West', '$28.0M', '$32.3M', '$35.9M', '$96.2M'],
            ['Total', '$85.3M', '$98.4M', '$102.7M', '$286.4M'],
        ]

        table = Table(data, colWidths=[1.2*inch, 1*inch, 1*inch, 1*inch, 1.2*inch])
        table.setStyle(TableStyle([
            ('BACKGROUND', (0, 0), (-1, 0), colors.grey),
            ('TEXTCOLOR', (0, 0), (-1, 0), colors.whitesmoke),
            ('ALIGN', (0, 0), (-1, -1), 'CENTER'),
            ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
            ('FONTSIZE', (0, 0), (-1, 0), 12),
            ('BOTTOMPADDING', (0, 0), (-1, 0), 12),
            ('BACKGROUND', (0, -1), (-1, -1), colors.lightgrey),
            ('FONTNAME', (0, -1), (-1, -1), 'Helvetica-Bold'),
            ('GRID', (0, 0), (-1, -1), 1, colors.black),
        ]))
        story.append(table)
        story.append(Spacer(1, 20))

        # Add decorative image
        decorative_path = TEST_DIR / "decorative.jpg"
        if decorative_path.exists():
            story.append(Paragraph("", styles['Normal']))  # Decorative divider
            img = Image(str(decorative_path), width=6*inch, height=1.5*inch)
            story.append(img)

        # Build PDF
        doc.build(story)
        print(f"✅ Created PDF with images: {output_path}")
        return output_path

    except ImportError as e:
        print(f"❌ reportlab not installed: {e}")
        return None


def main():
    """Create all test files."""
    print("=" * 60)
    print("Creating test files for cross-scanner integration testing")
    print("=" * 60)

    # Check for required test images
    required_images = ['decorative.jpg', 'infographic.jpg', 'bar_chart.png']
    missing = [img for img in required_images if not (TEST_DIR / img).exists()]
    if missing:
        print(f"\n⚠️  Missing test images: {missing}")
        print("Please download these images first or they will be skipped.")

    print("\n1. Creating PowerPoint with native charts...")
    create_pptx_with_charts()

    print("\n2. Creating PowerPoint with embedded images...")
    create_pptx_with_images()

    print("\n3. Creating PDF with images and tables...")
    create_pdf_with_images()

    print("\n" + "=" * 60)
    print("Test files created in:", TEST_DIR)
    print("=" * 60)


if __name__ == "__main__":
    main()
