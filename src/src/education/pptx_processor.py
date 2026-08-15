"""
PowerPoint Accessibility Scanner Module

This module provides functionality to:
1. Parse PPTX files and extract content
2. Detect missing alt text on images
3. Analyze text/background contrast ratios (WCAG 2.1 AA = 4.5:1)
4. Simulate color blindness and validate accessibility for CVD users (RGBlind integration)
5. Identify accessibility issues across all slides
6. Generate AI-powered remediation suggestions
7. Batch process entire directories
"""

from typing import List, Dict, Optional, Tuple, Any
from pydantic import BaseModel, computed_field
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
import os
import tempfile
from pathlib import Path

from src.utils.async_helpers import run_async_from_sync
from PIL import Image
from io import BytesIO


class ContrastIssue(BaseModel):
    """Contrast ratio issue"""

    slide_number: int
    shape_id: int
    shape_name: str
    foreground: str  # Hex color
    background: str  # Hex color
    contrast_ratio: float
    wcag_aa_pass: bool  # 4.5:1 for normal text
    wcag_aaa_pass: bool  # 7:1 for normal text
    suggested_fix: Optional[str] = None
    color_blindness_issues: Optional[List[Dict]] = None  # Issues for color-blind users


class AltTextIssue(BaseModel):
    """Missing or inaccurate alt text issue"""

    slide_number: int
    shape_id: int
    shape_name: str
    image_type: str
    has_alt_text: bool
    existing_alt_text: Optional[str] = None  # The current alt text if any
    suggested_alt_text: Optional[str] = None
    # Validation fields (when validate_alt_text is enabled)
    alt_text_validated: bool = False  # Whether AI validation was performed
    alt_text_accurate: Optional[bool] = None  # Whether existing alt text is accurate
    validation_issues: Optional[List[str]] = None  # Specific issues found
    validation_reasoning: Optional[str] = None  # AI explanation
    # Smart image analysis fields (cross-scanner integration)
    detected_image_type: Optional[str] = (
        None  # decorative, informative, functional, complex
    )
    is_decorative: bool = False  # True if image is decorative (needs empty alt)
    is_chart: bool = False  # True if detected as chart/graph/infographic
    detailed_description: Optional[str] = None  # For charts/complex images


class SlideTitleIssue(BaseModel):
    """Missing or problematic slide title issue (WCAG 1.3.1)"""

    slide_number: int
    issue_type: str  # missing_title, empty_title, duplicate_title
    existing_title: Optional[str] = None
    suggested_title: Optional[str] = None
    suggested_fix: str


class ImageOfTextIssue(BaseModel):
    """Image containing text issue (WCAG 1.4.5)"""

    slide_number: int
    shape_id: int
    shape_name: str
    detected_text: str  # The text detected in the image
    text_length: int  # Number of characters detected
    confidence: float  # OCR confidence score (0-100)
    suggested_fix: str


class AnimationIssue(BaseModel):
    """Animation accessibility issue (WCAG 2.2.2, 2.3.1)"""

    slide_number: int
    animation_index: int
    animation_type: str  # entrance, emphasis, exit, motion_path
    element_name: str  # Name of the animated element
    effect_name: Optional[str] = None  # Specific animation effect
    duration_ms: int  # Animation duration in milliseconds
    delay_ms: int = 0  # Delay before animation starts
    is_auto_start: bool = False  # True if animation starts automatically
    repeat_count: Optional[int] = None  # Number of times animation repeats
    issues: List[
        str
    ]  # List of specific issues: rapid_flash, auto_advance, motion_heavy
    suggested_fix: str


class EmbeddedMedia(BaseModel):
    """Embedded media in PPTX (WCAG 1.2.1, 1.2.2, 1.2.3)"""

    slide_number: int
    media_index: int
    media_type: str  # video, audio
    file_name: Optional[str] = None
    file_path: Optional[str] = None  # Path within PPTX
    duration_seconds: Optional[float] = None
    has_captions: bool = False
    has_transcript: bool = False  # Checked in speaker notes
    content_type: Optional[str] = None  # MIME type


class EmbeddedMediaIssue(BaseModel):
    """Embedded media accessibility issue"""

    slide_number: int
    media_index: int
    media_type: str  # video, audio
    file_name: Optional[str] = None
    issue_type: str  # missing_captions, missing_transcript, missing_audio_description
    recommendations: List[str]
    suggested_fix: str


class SlideAccessibilityIssues(BaseModel):
    """All accessibility issues for a single slide"""

    slide_number: int
    slide_title: Optional[str]
    contrast_issues: List[ContrastIssue]
    alt_text_issues: List[AltTextIssue]
    title_issues: List[SlideTitleIssue] = []  # Slide title issues (WCAG 1.3.1)
    image_of_text_issues: List[ImageOfTextIssue] = (
        []
    )  # Images containing text (WCAG 1.4.5)
    animation_issues: List[AnimationIssue] = []  # Animation issues
    embedded_media_issues: List[EmbeddedMediaIssue] = []  # Embedded media issues
    total_issues: int


class PowerPointProcessingResult(BaseModel):
    """Result of PowerPoint processing operation"""

    file_path: str
    file_name: str
    total_slides: int
    total_shapes: int
    total_images: int
    slides: List[SlideAccessibilityIssues]
    summary: Dict[
        str, int
    ]  # {"contrast_issues": X, "alt_text_issues": Y, "total_issues": Z}
    compliance_score: float  # 0-100
    remediation_suggestions: List[str]

    @computed_field
    @property
    def issues(self) -> List[Dict[str, Any]]:
        """Combined list of all issues for API compatibility.

        The demo routes expect a single 'issues' array, but PPTX processor
        stores issues in separate category arrays per slide. This computed
        field combines them all into a unified format.
        """
        all_issues: List[Dict[str, Any]] = []

        for slide in self.slides:
            slide_loc = f"Slide {slide.slide_number}"
            if slide.slide_title:
                slide_loc = f"{slide_loc} ({slide.slide_title})"

            # Contrast issues
            for issue in slide.contrast_issues:
                all_issues.append(
                    {
                        "id": f"contrast_{len(all_issues)}",
                        "category": "contrast",
                        "severity": "high" if issue.contrast_ratio < 3.0 else "medium",
                        "title": f"Low Contrast ({issue.contrast_ratio:.1f}:1)",
                        "description": f"Text contrast is {issue.contrast_ratio:.2f}:1, needs 4.5:1 for WCAG AA",
                        "location": f"{slide_loc}, Shape '{issue.shape_name}'",
                        "wcag_criterion": "WCAG 1.4.3",
                        "suggested_fix": issue.suggested_fix,
                    }
                )

            # Alt text issues
            for issue in slide.alt_text_issues:
                severity = "critical"
                if issue.is_decorative:
                    severity = "low"
                elif issue.has_alt_text and not issue.alt_text_accurate:
                    severity = "medium"

                if issue.is_decorative:
                    title = "Decorative Image — Needs Empty Alt Text"
                    desc = f"Image '{issue.shape_name}' is decorative and should have empty alt text (alt=\"\")"
                    fix = "Mark as decorative (set empty alt text)"
                elif not issue.has_alt_text:
                    title = "Missing Alt Text"
                    desc = f"Image '{issue.shape_name}' is missing alt text"
                    fix = issue.suggested_alt_text or "Add descriptive alt text"
                else:
                    title = "Inaccurate Alt Text"
                    desc = f"Image '{issue.shape_name}' has inaccurate alt text"
                    fix = issue.suggested_alt_text or "Update alt text"

                all_issues.append(
                    {
                        "id": f"alt_text_{len(all_issues)}",
                        "category": "alt_text",
                        "severity": severity,
                        "title": title,
                        "description": desc,
                        "location": f"{slide_loc}, Shape '{issue.shape_name}'",
                        "wcag_criterion": "WCAG 1.1.1",
                        "suggested_fix": fix,
                        "ai_generated": issue.suggested_alt_text is not None,
                        "generated_alt_text": issue.suggested_alt_text or None,
                        "is_decorative": issue.is_decorative,
                        "slide_index": slide.slide_number - 1,
                        "shape_name": issue.shape_name,
                    }
                )

            # Title issues
            for issue in slide.title_issues:
                if issue.issue_type == "missing_title":
                    desc = "Slide is missing a title, which is required for navigation"
                elif issue.issue_type == "empty_title":
                    desc = "Slide has an empty title placeholder"
                elif issue.issue_type == "duplicate_title":
                    desc = f"Slide title '{issue.existing_title}' is duplicated from another slide"
                else:
                    desc = f"Slide title issue: {issue.issue_type}"
                all_issues.append(
                    {
                        "id": f"title_{len(all_issues)}",
                        "category": "structure",
                        "severity": "high",
                        "title": issue.issue_type.replace("_", " ").title(),
                        "description": desc,
                        "location": slide_loc,
                        "wcag_criterion": "WCAG 1.3.1",
                        "suggested_fix": issue.suggested_fix,
                        "slide_index": slide.slide_number - 1,
                        "issue_type": issue.issue_type,
                    }
                )

            # Image of text issues
            for issue in slide.image_of_text_issues:
                all_issues.append(
                    {
                        "id": f"image_text_{len(all_issues)}",
                        "category": "image_of_text",
                        "severity": "medium",
                        "title": "Image Contains Text",
                        "description": f"Detected {issue.text_length} characters of text in image",
                        "location": f"{slide_loc}, Shape '{issue.shape_name}'",
                        "wcag_criterion": "WCAG 1.4.5",
                        "suggested_fix": issue.suggested_fix,
                    }
                )

        return all_issues


class PowerPointProcessor:
    """Process PowerPoint files for accessibility compliance"""

    def __init__(
        self,
        generate_alt_text: bool = False,
        validate_alt_text: bool = False,
        simulate_color_blindness: bool = True,
        detect_images_of_text: bool = False,
        progress_callback: callable = None,
    ):
        self.wcag_aa_ratio = 4.5  # WCAG 2.1 AA for normal text
        self.wcag_aaa_ratio = 7.0  # WCAG 2.1 AAA for normal text
        self.generate_alt_text = generate_alt_text
        self.validate_alt_text = (
            validate_alt_text  # Validate existing alt text accuracy
        )
        self.simulate_color_blindness = simulate_color_blindness
        self.detect_images_of_text = (
            detect_images_of_text  # OCR for images-of-text (WCAG 1.4.5)
        )
        self.progress_callback = progress_callback
        self.image_generator = None
        self.cvd_simulator = None
        self.tesseract_available = False

        # Lazy import to avoid circular dependencies
        if self.generate_alt_text or self.validate_alt_text:
            try:
                from .image_alt_text import ImageAltTextGenerator

                self.image_generator = ImageAltTextGenerator()
            except Exception as e:
                print(
                    f"[PowerPointProcessor] Warning: Could not initialize ImageAltTextGenerator: {e}"
                )
                self.generate_alt_text = False
                self.validate_alt_text = False

        # Initialize color blindness simulator if enabled
        if self.simulate_color_blindness:
            try:
                from .color_blindness_simulator import ColorBlindnessSimulator

                self.cvd_simulator = ColorBlindnessSimulator()
                print("[PowerPointProcessor] Color blindness simulation enabled")
            except Exception as e:
                print(
                    f"[PowerPointProcessor] Warning: Could not initialize ColorBlindnessSimulator: {e}"
                )
                self.simulate_color_blindness = False

        # Initialize tesseract for images-of-text detection if enabled
        if self.detect_images_of_text:
            try:
                import pytesseract

                # Test that tesseract is available
                pytesseract.get_tesseract_version()
                self.tesseract_available = True
                print(
                    "[PowerPointProcessor] Images-of-text detection enabled (pytesseract)"
                )
            except Exception as e:
                print(
                    f"[PowerPointProcessor] Warning: Could not initialize pytesseract: {e}. "
                    "Install tesseract-ocr to enable images-of-text detection."
                )
                self.detect_images_of_text = False

    def process_pptx(self, file_path: str) -> PowerPointProcessingResult:
        """
        Process a PowerPoint file and check accessibility

        Args:
            file_path: Path to PPTX file

        Returns:
            PowerPointProcessingResult with all accessibility issues
        """
        prs = Presentation(file_path)
        file_name = os.path.basename(file_path)
        total_slides = len(prs.slides)

        # Report initial progress
        if self.progress_callback:
            self.progress_callback(0, total_slides, "Loading presentation...")

        # Extract presentation-level context for AI fixes
        presentation_context = self._extract_presentation_context(prs)

        total_shapes = 0
        total_images = 0
        slides_issues = []

        # Process each slide
        for slide_idx, slide in enumerate(prs.slides, start=1):
            # Report progress per slide
            if self.progress_callback:
                self.progress_callback(
                    slide_idx,
                    total_slides,
                    f"Analyzing slide {slide_idx} of {total_slides}...",
                )
            slide_title = self._get_slide_title(slide)
            slide_context = self._extract_slide_context(
                slide, slide_idx, presentation_context
            )
            contrast_issues = []
            alt_text_issues = []
            title_issues = []
            image_of_text_issues = []
            animation_issues = []

            # Analyze slide animations (WCAG 2.2.2, 2.3.1)
            animation_issues = self._analyze_animations(slide, slide_idx, file_path)

            # Check embedded media (WCAG 1.2.1, 1.2.2, 1.2.3)
            embedded_media_issues = self._check_embedded_media(
                slide, slide_idx, file_path
            )

            # Check slide title (WCAG 1.3.1 - Info and Relationships)
            title_issue = self._check_slide_title(
                slide_idx, slide, slide_title, presentation_context
            )
            if title_issue:
                title_issues.append(title_issue)

            # Analyze all shapes on the slide
            for shape in slide.shapes:
                total_shapes += 1

                # Check for images
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    total_images += 1
                    alt_issue = self._check_alt_text(slide_idx, shape, slide_context)
                    if alt_issue:
                        alt_text_issues.append(alt_issue)

                    # Check for images-of-text (WCAG 1.4.5)
                    if self.detect_images_of_text:
                        text_issue = self._check_image_of_text(slide_idx, shape)
                        if text_issue:
                            image_of_text_issues.append(text_issue)

                # Check text contrast
                if hasattr(shape, "text_frame"):
                    contrast_issue = self._check_text_contrast(
                        slide_idx, shape, slide_context
                    )
                    if contrast_issue:
                        contrast_issues.append(contrast_issue)

            # Create slide summary
            slide_issues = SlideAccessibilityIssues(
                slide_number=slide_idx,
                slide_title=slide_title,
                contrast_issues=contrast_issues,
                alt_text_issues=alt_text_issues,
                title_issues=title_issues,
                image_of_text_issues=image_of_text_issues,
                animation_issues=animation_issues,
                embedded_media_issues=embedded_media_issues,
                total_issues=len(contrast_issues)
                + len(alt_text_issues)
                + len(title_issues)
                + len(image_of_text_issues)
                + len(animation_issues)
                + len(embedded_media_issues),
            )
            slides_issues.append(slide_issues)

        # Calculate summary and compliance score
        summary = self._calculate_summary(slides_issues)
        compliance_score = self._calculate_compliance_score(
            summary, total_shapes, total_images
        )
        remediation_suggestions = self._generate_remediation_suggestions(summary)

        return PowerPointProcessingResult(
            file_path=file_path,
            file_name=file_name,
            total_slides=len(prs.slides),
            total_shapes=total_shapes,
            total_images=total_images,
            slides=slides_issues,
            summary=summary,
            compliance_score=compliance_score,
            remediation_suggestions=remediation_suggestions,
        )

    def _get_slide_title(self, slide) -> Optional[str]:
        """Extract title from slide if present"""
        try:
            if slide.shapes.title:
                return slide.shapes.title.text
        except Exception:
            pass
        return None

    def _check_slide_title(
        self,
        slide_number: int,
        slide,
        slide_title: Optional[str],
        presentation_context: Dict,
    ) -> Optional[SlideTitleIssue]:
        """
        Check slide title for accessibility issues (WCAG 1.3.1).

        Detects:
        - Missing slide titles (no title placeholder or title shape)
        - Empty slide titles (title placeholder exists but has no text)
        - Duplicate slide titles (same title used on multiple slides)

        Args:
            slide_number: The 1-based slide number
            slide: The slide object
            slide_title: The extracted title text (or None)
            presentation_context: Context including all slide titles

        Returns:
            SlideTitleIssue if an issue is found, None otherwise
        """
        # Check for missing or empty title
        if not slide_title or not slide_title.strip():
            # Try to generate a suggested title from slide content
            suggested_title = self._generate_slide_title_suggestion(slide, slide_number)

            # Determine if title placeholder exists but is empty vs. completely missing
            has_title_placeholder = False
            try:
                has_title_placeholder = slide.shapes.title is not None
            except Exception:
                pass

            if has_title_placeholder and (not slide_title or not slide_title.strip()):
                return SlideTitleIssue(
                    slide_number=slide_number,
                    issue_type="empty_title",
                    existing_title=slide_title,
                    suggested_title=suggested_title,
                    suggested_fix=f'Add a descriptive title to slide {slide_number}. Suggested: "{suggested_title}"',
                )
            else:
                return SlideTitleIssue(
                    slide_number=slide_number,
                    issue_type="missing_title",
                    existing_title=None,
                    suggested_title=suggested_title,
                    suggested_fix=f'Add a title placeholder with descriptive text to slide {slide_number}. Suggested: "{suggested_title}"',
                )

        # Check for duplicate titles (same title used elsewhere)
        slide_titles = presentation_context.get("slide_titles", [])
        title_normalized = slide_title.strip().lower()

        # Count occurrences of this title
        occurrences = sum(
            1
            for st in slide_titles
            if st.get("title", "").strip().lower() == title_normalized
        )

        if occurrences > 1:
            return SlideTitleIssue(
                slide_number=slide_number,
                issue_type="duplicate_title",
                existing_title=slide_title,
                suggested_title=f"{slide_title} ({slide_number})",
                suggested_fix=f'Slide {slide_number} has a duplicate title "{slide_title}". Make each slide title unique for better navigation.',
            )

        return None

    def _generate_slide_title_suggestion(self, slide, slide_number: int) -> str:
        """
        Generate a suggested title for a slide based on its content.

        Args:
            slide: The slide object
            slide_number: The 1-based slide number

        Returns:
            A suggested title string
        """
        # Try to extract meaningful text from the slide
        text_content = []

        for shape in slide.shapes:
            # Skip title placeholder (we know it's empty/missing)
            try:
                if slide.shapes.title and shape == slide.shapes.title:
                    continue
            except Exception:
                pass

            # Extract text from text frames
            if hasattr(shape, "text_frame"):
                text = shape.text_frame.text.strip()
                if text:
                    # Get first line or first 50 characters
                    first_line = text.split("\n")[0].strip()
                    if first_line and len(first_line) > 3:
                        text_content.append(first_line[:80])

        # Use the most prominent text as title suggestion
        if text_content:
            # Pick the first substantial piece of text
            for text in text_content:
                if len(text) >= 5:
                    return text[:60]

        # Fallback to generic title
        return f"Slide {slide_number}"

    def _extract_presentation_context(self, prs: Presentation) -> Dict:
        """
        Extract presentation-level context for AI to understand the overall structure.
        This helps AI generate more relevant alt text and fixes.
        """
        context = {
            "total_slides": len(prs.slides),
            "slide_titles": [],
            "theme": None,  # Could be extracted from prs if needed
        }

        for slide_idx, slide in enumerate(prs.slides, start=1):
            title = self._get_slide_title(slide)
            context["slide_titles"].append(
                {
                    "slide_number": slide_idx,
                    "title": title or f"[Untitled Slide {slide_idx}]",
                }
            )

        return context

    def _extract_slide_context(
        self, slide, slide_idx: int, presentation_context: Dict
    ) -> Dict:
        """
        Extract detailed context for a specific slide.
        Includes text content, structure, and position within presentation.
        """
        context = {
            "slide_number": slide_idx,
            "slide_title": self._get_slide_title(slide),
            "total_slides": presentation_context.get("total_slides", 1),
            "text_content": [],
            "bullet_points": [],
            "image_count": 0,
            "table_count": 0,
            "previous_slide_title": None,
            "next_slide_title": None,
        }

        # Get adjacent slide titles for context
        slide_titles = presentation_context.get("slide_titles", [])
        if slide_idx > 1 and len(slide_titles) >= slide_idx - 1:
            context["previous_slide_title"] = slide_titles[slide_idx - 2].get("title")
        if slide_idx < len(slide_titles):
            context["next_slide_title"] = (
                slide_titles[slide_idx].get("title")
                if slide_idx < len(slide_titles)
                else None
            )

        # Extract text content from all shapes
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                context["image_count"] += 1

            if hasattr(shape, "text_frame"):
                text = shape.text_frame.text.strip()
                if text:
                    # Check if it's bullet points
                    if "\n" in text:
                        for line in text.split("\n"):
                            line = line.strip()
                            if line:
                                context["bullet_points"].append(line[:100])
                    else:
                        context["text_content"].append(text[:200])

            # Check for tables - need try/except because hasattr returns True
            # for chart shapes that inherit from GraphicFrame but accessing
            # .table raises ValueError
            try:
                if hasattr(shape, "table") and shape.table:
                    context["table_count"] += 1
            except ValueError:
                # This is a chart, not a table - ignore
                pass

        return context

    def _build_context_string_for_alt_text(self, slide_context: Dict) -> str:
        """Build a context string to help AI generate relevant alt text"""
        parts = []

        parts.append(
            f"Slide {slide_context['slide_number']} of {slide_context['total_slides']}"
        )

        if slide_context.get("slide_title"):
            parts.append(f"Slide Title: \"{slide_context['slide_title']}\"")

        if slide_context.get("previous_slide_title"):
            parts.append(f"Previous Slide: \"{slide_context['previous_slide_title']}\"")

        if slide_context.get("text_content"):
            # Include relevant text content (limit to first 3 items)
            text_items = slide_context["text_content"][:3]
            parts.append("Text on slide:")
            for text in text_items:
                parts.append(f'  - "{text}"')

        if slide_context.get("bullet_points"):
            # Include some bullet points for context
            bullets = slide_context["bullet_points"][:5]
            parts.append("Key points on slide:")
            for bullet in bullets:
                parts.append(f"  • {bullet}")

        if slide_context.get("image_count", 0) > 1:
            parts.append(
                f"Note: This slide has {slide_context['image_count']} images total"
            )

        return "\n".join(parts)

    def _extract_image_from_shape(self, shape) -> Optional[str]:
        """
        Extract image from PowerPoint shape and save to temporary file

        Returns:
            Path to temporary image file, or None if extraction failed
        """
        try:
            # Get image bytes from shape
            image_blob = shape.image.blob
            image = Image.open(BytesIO(image_blob))

            # Save to temporary file
            with tempfile.NamedTemporaryFile(delete=False, suffix=".png") as tmp:
                image.save(tmp, format="PNG")
                return tmp.name
        except Exception as e:
            print(f"[PowerPointProcessor] Failed to extract image: {e}")
            return None

    def _check_image_of_text(
        self, slide_number: int, shape
    ) -> Optional[ImageOfTextIssue]:
        """
        Check if an image contains text using OCR (WCAG 1.4.5 - Images of Text).

        WCAG 1.4.5 requires that images of text should only be used for:
        - Pure decoration
        - Essential presentation (e.g., logos)

        Text should be rendered as actual text, not as images.

        Args:
            slide_number: The 1-based slide number
            shape: The image shape to check

        Returns:
            ImageOfTextIssue if significant text is detected, None otherwise
        """
        if not self.tesseract_available:
            return None

        # Minimum text threshold to consider as "image of text"
        # Small amounts of text (like labels) are often acceptable
        MIN_TEXT_LENGTH = 20  # Characters
        MIN_CONFIDENCE = 40  # OCR confidence threshold

        try:
            import pytesseract

            # Extract image to temp file
            image_path = self._extract_image_from_shape(shape)
            if not image_path:
                return None

            try:
                # Open and preprocess image for better OCR
                image = Image.open(image_path)

                # Convert to RGB if necessary (grayscale can help OCR)
                if image.mode != "RGB":
                    image = image.convert("RGB")

                # Get OCR data with confidence scores
                ocr_data = pytesseract.image_to_data(
                    image,
                    output_type=pytesseract.Output.DICT,
                    config="--oem 3 --psm 6",  # OCR Engine Mode 3, Page Seg Mode 6
                )

                # Filter and combine text with confidence
                detected_words = []
                confidence_scores = []

                for i, word in enumerate(ocr_data["text"]):
                    word = word.strip()
                    conf = int(ocr_data["conf"][i])

                    # Only include words with decent confidence
                    if word and conf > MIN_CONFIDENCE:
                        detected_words.append(word)
                        confidence_scores.append(conf)

                # Combine detected text
                detected_text = " ".join(detected_words)
                avg_confidence = (
                    sum(confidence_scores) / len(confidence_scores)
                    if confidence_scores
                    else 0
                )

                # Clean up temp file
                try:
                    os.unlink(image_path)
                except Exception:
                    pass

                # Check if significant text was found
                if len(detected_text) >= MIN_TEXT_LENGTH:
                    # Truncate text preview to reasonable length
                    text_preview = (
                        detected_text[:150] + "..."
                        if len(detected_text) > 150
                        else detected_text
                    )

                    return ImageOfTextIssue(
                        slide_number=slide_number,
                        shape_id=shape.shape_id,
                        shape_name=shape.name,
                        detected_text=text_preview,
                        text_length=len(detected_text),
                        confidence=round(avg_confidence, 1),
                        suggested_fix=(
                            f"This image contains text ({len(detected_text)} characters). "
                            "Consider using real text instead of an image for accessibility and "
                            "searchability. Exception: decorative text or essential logos. (WCAG 1.4.5)"
                        ),
                    )

            except Exception as e:
                print(f"[PowerPointProcessor] OCR failed for image: {e}")
                # Clean up temp file on error
                try:
                    os.unlink(image_path)
                except Exception:
                    pass

        except Exception as e:
            print(f"[PowerPointProcessor] Image-of-text detection failed: {e}")

        return None

    def _analyze_animations(
        self, slide, slide_number: int, file_path: str
    ) -> List[AnimationIssue]:
        """
        Analyze slide animations for accessibility issues.

        Checks for:
        - Rapid animations that could cause seizures (WCAG 2.3.1)
        - Auto-starting animations without user control (WCAG 2.2.2)
        - Excessive motion that could cause vestibular issues
        - Very short animations that flash content

        Args:
            slide: python-pptx slide object
            slide_number: 1-indexed slide number
            file_path: Path to PPTX file for direct XML access

        Returns:
            List of AnimationIssue objects
        """
        import zipfile
        from xml.etree import ElementTree as ET

        issues = []

        # PPTX animation namespaces (for future XPath queries)
        _namespaces = {  # noqa: F841
            "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
            "p": "http://schemas.openxmlformats.org/presentationml/2006/main",
            "r": "http://schemas.openxmlformats.org/officeDocument/2006/relationships",
        }

        try:
            with zipfile.ZipFile(file_path, "r") as pptx_zip:
                # Find the slide XML file
                slide_xml_path = f"ppt/slides/slide{slide_number}.xml"

                if slide_xml_path not in pptx_zip.namelist():
                    return issues

                with pptx_zip.open(slide_xml_path) as f:
                    tree = ET.parse(f)
                    root = tree.getroot()

                    # Find timing/animation elements
                    # Animations are typically in p:timing/p:tnLst
                    anim_index = 0

                    for timing in root.iter():
                        tag_name = (
                            timing.tag.split("}")[-1]
                            if "}" in timing.tag
                            else timing.tag
                        )

                        # Check for animation effect elements
                        if tag_name in [
                            "anim",
                            "animEffect",
                            "set",
                            "animMotion",
                            "animClr",
                            "animScale",
                            "animRot",
                        ]:
                            anim_issues = []
                            anim_type = self._get_animation_type(tag_name)
                            duration_ms = self._get_animation_duration(timing)
                            is_auto = self._is_auto_start(timing)
                            repeat = self._get_repeat_count(timing)
                            effect_name = (
                                timing.get("filter") or timing.get("prst") or tag_name
                            )

                            # Get target element name
                            element_name = self._get_animation_target_name(
                                timing, slide
                            )

                            # Check for rapid flash risk (duration < 50ms with repeat)
                            if (
                                duration_ms > 0
                                and duration_ms < 50
                                and repeat
                                and repeat > 3
                            ):
                                anim_issues.append("rapid_flash")

                            # Check for auto-start without user control
                            if is_auto:
                                anim_issues.append("auto_advance")

                            # Check for excessive motion (motion paths)
                            if tag_name == "animMotion":
                                anim_issues.append("motion_heavy")

                            # Check for very rapid animations (< 100ms)
                            if 0 < duration_ms < 100:
                                anim_issues.append("rapid_animation")

                            # Check for high repeat count
                            if repeat and repeat > 5:
                                anim_issues.append("excessive_repeat")

                            if anim_issues:
                                suggested_fix = self._generate_animation_fix(
                                    anim_issues, anim_type
                                )

                                issues.append(
                                    AnimationIssue(
                                        slide_number=slide_number,
                                        animation_index=anim_index,
                                        animation_type=anim_type,
                                        element_name=element_name
                                        or f"Element {anim_index + 1}",
                                        effect_name=effect_name,
                                        duration_ms=duration_ms,
                                        is_auto_start=is_auto,
                                        repeat_count=repeat,
                                        issues=anim_issues,
                                        suggested_fix=suggested_fix,
                                    )
                                )

                            anim_index += 1

                    # Also check for slide transition timing
                    for trans in root.iter():
                        tag_name = (
                            trans.tag.split("}")[-1] if "}" in trans.tag else trans.tag
                        )

                        if tag_name == "transition":
                            # Check for auto-advance slide transition
                            adv_tm = trans.get("advTm")  # Auto-advance time in ms
                            if adv_tm:
                                try:
                                    adv_ms = int(adv_tm)
                                    if (
                                        adv_ms < 5000
                                    ):  # Auto-advance in less than 5 seconds
                                        issues.append(
                                            AnimationIssue(
                                                slide_number=slide_number,
                                                animation_index=anim_index,
                                                animation_type="transition",
                                                element_name="Slide Transition",
                                                effect_name="Auto-advance",
                                                duration_ms=adv_ms,
                                                is_auto_start=True,
                                                issues=[
                                                    "auto_advance",
                                                    "rapid_transition",
                                                ],
                                                suggested_fix="Auto-advancing slides can be disorienting. "
                                                "Consider requiring user click to advance, or "
                                                "increase time to at least 5 seconds. (WCAG 2.2.2)",
                                            )
                                        )
                                        anim_index += 1
                                except ValueError:
                                    pass

        except Exception as e:
            print(f"[PowerPointProcessor] Animation analysis failed: {e}")

        return issues

    def _get_animation_type(self, tag_name: str) -> str:
        """Map XML tag to animation type category."""
        type_map = {
            "anim": "entrance",
            "animEffect": "emphasis",
            "set": "emphasis",
            "animMotion": "motion_path",
            "animClr": "emphasis",
            "animScale": "emphasis",
            "animRot": "emphasis",
        }
        return type_map.get(tag_name, "other")

    def _get_animation_duration(self, element) -> int:
        """Extract animation duration in milliseconds."""
        # Check common duration attributes
        dur = element.get("dur")
        if dur:
            if dur == "indefinite":
                return 0
            try:
                # Duration might be in various formats
                if dur.isdigit():
                    return int(dur)
            except (ValueError, AttributeError):
                pass

        # Check child cTn (common time node) element
        for child in element:
            child_tag = child.tag.split("}")[-1] if "}" in child.tag else child.tag
            if child_tag == "cTn":
                dur = child.get("dur")
                if dur and dur.isdigit():
                    return int(dur)

        return 500  # Default animation duration

    def _is_auto_start(self, element) -> bool:
        """Check if animation starts automatically."""
        # Check for nodeType or stCondLst indicating auto-start
        for child in element:
            child_tag = child.tag.split("}")[-1] if "}" in child.tag else child.tag
            if child_tag == "cTn":
                node_type = child.get("nodeType")
                if node_type in ["afterPrevious", "withPrevious"]:
                    return True

        return False

    def _get_repeat_count(self, element) -> Optional[int]:
        """Get animation repeat count."""
        for child in element:
            child_tag = child.tag.split("}")[-1] if "}" in child.tag else child.tag
            if child_tag == "cTn":
                repeat = child.get("repeatCount")
                if repeat:
                    if repeat == "indefinite":
                        return 1000  # Treat indefinite as high number
                    try:
                        return int(repeat)
                    except ValueError:
                        pass
        return None

    def _get_animation_target_name(self, element, slide) -> Optional[str]:
        """Get name of the element being animated."""
        # Try to find target reference
        for child in element:
            child_tag = child.tag.split("}")[-1] if "}" in child.tag else child.tag
            if child_tag == "tgtEl":
                for target in child:
                    target_tag = (
                        target.tag.split("}")[-1] if "}" in target.tag else target.tag
                    )
                    if target_tag == "spTgt":
                        shape_id = target.get("spid")
                        if shape_id:
                            # Try to find shape name from slide
                            try:
                                for shape in slide.shapes:
                                    if str(shape.shape_id) == shape_id:
                                        return shape.name
                            except Exception:
                                pass
                            return f"Shape {shape_id}"
        return None

    def _generate_animation_fix(self, issues: List[str], anim_type: str) -> str:
        """Generate suggested fix based on animation issues."""
        fixes = []

        if "rapid_flash" in issues:
            fixes.append(
                "This animation may flash rapidly, risking seizures for photosensitive users. "
                "Reduce repeat count or increase duration to > 50ms per cycle."
            )

        if "auto_advance" in issues:
            fixes.append(
                "Animation starts automatically without user control. "
                "Consider requiring a click to start, or provide pause/stop controls."
            )

        if "motion_heavy" in issues:
            fixes.append(
                "Motion path animation may cause vestibular discomfort. "
                "Consider reducing motion or providing a reduced-motion alternative."
            )

        if "rapid_animation" in issues:
            fixes.append(
                "Animation is very rapid (< 100ms). "
                "Consider slowing down for better accessibility."
            )

        if "excessive_repeat" in issues:
            fixes.append(
                "Animation repeats many times. "
                "Consider reducing repetitions or allowing users to stop it."
            )

        if fixes:
            return " ".join(fixes) + " (WCAG 2.2.2, 2.3.1)"

        return f"Review {anim_type} animation for accessibility concerns."

    def _check_embedded_media(
        self, slide, slide_number: int, file_path: str
    ) -> List[EmbeddedMediaIssue]:
        """
        Check embedded video/audio in slide for accessibility.

        Verifies:
        - Captions/subtitles for video content (WCAG 1.2.2)
        - Transcripts available (WCAG 1.2.1, 1.2.3)
        - Audio descriptions for video (WCAG 1.2.3, 1.2.5)

        Args:
            slide: python-pptx slide object
            slide_number: 1-indexed slide number
            file_path: Path to PPTX file

        Returns:
            List of EmbeddedMediaIssue objects
        """
        import zipfile
        from xml.etree import ElementTree as ET

        issues = []
        media_found = []

        # PPTX namespaces (reserved for future XPath queries)
        _namespaces = {  # noqa: F841
            "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
            "p": "http://schemas.openxmlformats.org/presentationml/2006/main",
            "r": "http://schemas.openxmlformats.org/officeDocument/2006/relationships",
        }

        try:
            with zipfile.ZipFile(file_path, "r") as pptx_zip:
                # Check for media files in ppt/media/ (for future media inventory)
                _media_files = [  # noqa: F841
                    f for f in pptx_zip.namelist() if f.startswith("ppt/media/")
                ]

                # Identify video and audio files (for future type filtering)
                _video_extensions = {
                    ".mp4",
                    ".avi",
                    ".mov",
                    ".wmv",
                    ".m4v",
                    ".webm",
                }  # noqa: F841
                _audio_extensions = {
                    ".mp3",
                    ".wav",
                    ".m4a",
                    ".wma",
                    ".ogg",
                    ".aac",
                }  # noqa: F841

                # Parse slide relationships to find media references
                slide_rels_path = f"ppt/slides/_rels/slide{slide_number}.xml.rels"

                slide_media = []
                if slide_rels_path in pptx_zip.namelist():
                    with pptx_zip.open(slide_rels_path) as f:
                        rels_tree = ET.parse(f)
                        rels_root = rels_tree.getroot()

                        for rel in rels_root:
                            target = rel.get("Target", "")
                            rel_type = rel.get("Type", "")

                            # Check for video relationship
                            if (
                                "video" in rel_type.lower()
                                or "media" in rel_type.lower()
                            ):
                                slide_media.append(
                                    {
                                        "target": target,
                                        "type": (
                                            "video"
                                            if "video" in rel_type.lower()
                                            else "media"
                                        ),
                                    }
                                )

                            # Check for audio relationship
                            if "audio" in rel_type.lower():
                                slide_media.append(
                                    {
                                        "target": target,
                                        "type": "audio",
                                    }
                                )

                # Also check for media directly in slide XML
                slide_xml_path = f"ppt/slides/slide{slide_number}.xml"
                if slide_xml_path in pptx_zip.namelist():
                    with pptx_zip.open(slide_xml_path) as f:
                        slide_tree = ET.parse(f)
                        slide_root = slide_tree.getroot()

                        # Look for video/audio elements
                        for elem in slide_root.iter():
                            tag_name = (
                                elem.tag.split("}")[-1] if "}" in elem.tag else elem.tag
                            )

                            if tag_name in ["video", "audio", "videoFile", "audioFile"]:
                                embed_ref = elem.get(
                                    "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}embed"
                                )
                                link_ref = elem.get(
                                    "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}link"
                                )

                                media_type = (
                                    "video" if "video" in tag_name.lower() else "audio"
                                )
                                slide_media.append(
                                    {
                                        "target": embed_ref or link_ref or "unknown",
                                        "type": media_type,
                                    }
                                )

                # Check speaker notes for transcripts
                has_transcript_in_notes = self._check_speaker_notes_for_transcript(
                    slide, slide_number, pptx_zip
                )

                # Process found media
                for media_idx, media_info in enumerate(slide_media):
                    media_type = media_info["type"]
                    target = media_info["target"]
                    file_name = os.path.basename(target) if target else None

                    # Check for captions (look for .vtt, .srt files)
                    has_captions = self._check_for_captions(pptx_zip, target)

                    media_item = EmbeddedMedia(
                        slide_number=slide_number,
                        media_index=media_idx,
                        media_type=media_type,
                        file_name=file_name,
                        file_path=target,
                        has_captions=has_captions,
                        has_transcript=has_transcript_in_notes,
                    )
                    media_found.append(media_item)

                    # Generate issues for missing accessibility features
                    issue_recommendations = []
                    issue_type = None

                    if media_type == "video":
                        if not has_captions:
                            issue_type = "missing_captions"
                            issue_recommendations.append(
                                "Add synchronized captions or subtitles to the video (WCAG 1.2.2)"
                            )
                            issue_recommendations.append(
                                "Consider using a video hosting platform that supports automatic captions"
                            )

                        if not has_transcript_in_notes:
                            if not issue_type:
                                issue_type = "missing_transcript"
                            issue_recommendations.append(
                                "Provide a text transcript in the speaker notes or nearby slide (WCAG 1.2.1)"
                            )

                        # Audio descriptions recommendation for video
                        issue_recommendations.append(
                            "Consider adding audio descriptions for visual content "
                            "not conveyed by the soundtrack (WCAG 1.2.3, 1.2.5)"
                        )

                    elif media_type == "audio":
                        if not has_transcript_in_notes:
                            issue_type = "missing_transcript"
                            issue_recommendations.append(
                                "Provide a text transcript for the audio content (WCAG 1.2.1)"
                            )
                            issue_recommendations.append(
                                "Add transcript to speaker notes or a nearby slide"
                            )

                    if issue_recommendations:
                        issues.append(
                            EmbeddedMediaIssue(
                                slide_number=slide_number,
                                media_index=media_idx,
                                media_type=media_type,
                                file_name=file_name,
                                issue_type=issue_type or "accessibility_review",
                                recommendations=issue_recommendations,
                                suggested_fix=(
                                    issue_recommendations[0]
                                    if issue_recommendations
                                    else ""
                                ),
                            )
                        )

        except Exception as e:
            print(f"[PowerPointProcessor] Embedded media check failed: {e}")

        return issues

    def _check_speaker_notes_for_transcript(
        self, slide, slide_number: int, pptx_zip
    ) -> bool:
        """
        Check if speaker notes contain a transcript.

        Looks for transcript indicators in speaker notes:
        - "Transcript:" or "TRANSCRIPT:"
        - Long text that might be a transcript (>200 chars)
        - [Transcript] markers

        Args:
            slide: python-pptx slide object
            slide_number: 1-indexed slide number
            pptx_zip: ZipFile object for PPTX

        Returns:
            True if transcript appears to be present
        """
        try:
            # Check python-pptx notes
            if slide.has_notes_slide:
                notes_slide = slide.notes_slide
                notes_text = (
                    notes_slide.notes_text_frame.text
                    if notes_slide.notes_text_frame
                    else ""
                )

                if notes_text:
                    notes_lower = notes_text.lower()

                    # Check for transcript markers
                    if "transcript" in notes_lower:
                        return True

                    # Check for substantial text that could be transcript
                    if len(notes_text) > 200:
                        return True

            # Also check notes XML directly
            notes_xml_path = f"ppt/notesSlides/notesSlide{slide_number}.xml"
            if notes_xml_path in pptx_zip.namelist():
                from xml.etree import ElementTree as ET

                with pptx_zip.open(notes_xml_path) as f:
                    notes_tree = ET.parse(f)
                    notes_root = notes_tree.getroot()

                    # Extract all text
                    text_content = []
                    for elem in notes_root.iter():
                        if elem.text:
                            text_content.append(elem.text)

                    full_text = " ".join(text_content)
                    if "transcript" in full_text.lower() or len(full_text) > 200:
                        return True

        except Exception as e:
            print(f"[PowerPointProcessor] Notes check failed: {e}")

        return False

    def _check_for_captions(self, pptx_zip, media_target: str) -> bool:
        """
        Check if captions/subtitles exist for media file.

        Looks for:
        - .vtt (WebVTT) files
        - .srt (SubRip) files
        - Caption tracks in media metadata

        Args:
            pptx_zip: ZipFile object for PPTX
            media_target: Target path of media file

        Returns:
            True if captions appear to be present
        """
        if not media_target:
            return False

        try:
            # Get base name without extension
            base_name = os.path.splitext(os.path.basename(media_target))[0]

            # Look for caption files with same base name
            caption_extensions = [".vtt", ".srt", ".sbv", ".sub", ".ass"]

            for file_name in pptx_zip.namelist():
                file_base = os.path.splitext(os.path.basename(file_name))[0]
                file_ext = os.path.splitext(file_name)[1].lower()

                # Check if caption file matches media file
                if file_ext in caption_extensions:
                    if base_name.lower() in file_base.lower():
                        return True

            # Check for embedded captions in media folder
            if "ppt/media/" in media_target:
                for file_name in pptx_zip.namelist():
                    if file_name.startswith("ppt/media/") and any(
                        file_name.endswith(ext) for ext in caption_extensions
                    ):
                        return True

        except Exception as e:
            print(f"[PowerPointProcessor] Caption check failed: {e}")

        return False

    def _check_alt_text(
        self, slide_number: int, shape, slide_context: Dict = None
    ) -> Optional[AltTextIssue]:
        """
        Check if image has alt text, generate if missing, and validate if existing.

        SMART IMAGE ANALYSIS (Cross-Scanner Integration):
        1. First detect image type (decorative/informative/functional/complex)
        2. If decorative: report that empty alt="" is recommended
        3. If complex (chart/graph/infographic): use describe_chart_or_graph()
        4. Otherwise: use standard generate_alt_text()
        """
        try:
            alt_text = shape._element.get("descr", "")
            has_alt = alt_text and alt_text.strip() != ""

            if not has_alt:
                # Image is MISSING alt text - use smart analysis
                suggested_alt = None
                detected_type = None
                is_decorative = False
                is_chart = False
                detailed_desc = None

                # Generate alt text with AI if enabled (with smart type detection)
                if self.generate_alt_text and self.image_generator:
                    try:
                        image_path = self._extract_image_from_shape(shape)

                        if image_path:
                            if slide_context:
                                context = self._build_context_string_for_alt_text(
                                    slide_context
                                )
                            else:
                                context = f"Image from slide {slide_number}"

                            # PHASE 1: Detect image type first
                            type_result = run_async_from_sync(
                                self.image_generator.detect_image_type(
                                    image_path=image_path, context=context
                                )
                            )

                            if type_result.get("success"):
                                detected_type = type_result.get(
                                    "image_type", "informative"
                                )
                                is_decorative = type_result.get("is_decorative", False)
                                image_purpose = type_result.get("image_purpose", "")

                                # Check if it's a chart/graph/infographic
                                is_chart = detected_type == "complex" or any(
                                    term in image_purpose.lower()
                                    for term in [
                                        "chart",
                                        "graph",
                                        "plot",
                                        "diagram",
                                        "infographic",
                                        "data visualization",
                                    ]
                                )

                                print(
                                    f"[PowerPointProcessor] Slide {slide_number}: Detected image type = {detected_type}, is_chart = {is_chart}, is_decorative = {is_decorative}"
                                )

                                # PHASE 2: Generate appropriate description based on type
                                if is_decorative:
                                    # Decorative images need empty alt text
                                    suggested_alt = ""  # WCAG compliant empty alt
                                    print(
                                        f'[PowerPointProcessor] Slide {slide_number}: Image is DECORATIVE - recommending empty alt=""'
                                    )
                                elif is_chart:
                                    # Charts/graphs need detailed descriptions
                                    chart_result = run_async_from_sync(
                                        self.image_generator.describe_chart_or_graph(
                                            image_path=image_path,
                                            context=context,
                                            detail_level="standard",
                                        )
                                    )
                                    if chart_result.get("success"):
                                        suggested_alt = chart_result.get(
                                            "short_description", ""
                                        )
                                        detailed_desc = chart_result.get(
                                            "detailed_description", ""
                                        )
                                        chart_type = chart_result.get(
                                            "chart_type", "chart"
                                        )
                                        print(
                                            f"[PowerPointProcessor] Slide {slide_number}: Generated CHART description ({chart_type}): {suggested_alt[:50]}..."
                                        )
                                    else:
                                        # Fallback to regular alt text
                                        result = run_async_from_sync(
                                            self.image_generator.generate_alt_text(
                                                image_path=image_path,
                                                context=context,
                                                educational_context=True,
                                            )
                                        )
                                        if result.get("success"):
                                            suggested_alt = result.get("alt_text")
                                else:
                                    # Regular informative images
                                    result = run_async_from_sync(
                                        self.image_generator.generate_alt_text(
                                            image_path=image_path,
                                            context=context,
                                            educational_context=True,
                                        )
                                    )
                                    if result.get("success"):
                                        suggested_alt = result.get("alt_text")
                                        print(
                                            f"[PowerPointProcessor] Slide {slide_number}: Generated standard alt text: {suggested_alt[:50]}..."
                                        )
                            else:
                                # Type detection failed, fallback to standard alt text generation
                                result = run_async_from_sync(
                                    self.image_generator.generate_alt_text(
                                        image_path=image_path,
                                        context=context,
                                        educational_context=True,
                                    )
                                )
                                if result.get("success"):
                                    suggested_alt = result.get("alt_text")
                                    print(
                                        f"[PowerPointProcessor] Slide {slide_number}: Generated fallback alt text: {suggested_alt[:50]}..."
                                    )

                            try:
                                os.unlink(image_path)
                            except Exception:
                                pass
                    except Exception as e:
                        print(f"[PowerPointProcessor] Alt text generation failed: {e}")

                return AltTextIssue(
                    slide_number=slide_number,
                    shape_id=shape.shape_id,
                    shape_name=shape.name,
                    image_type="picture",
                    has_alt_text=False,
                    existing_alt_text=None,
                    suggested_alt_text=suggested_alt,
                    detected_image_type=detected_type,
                    is_decorative=is_decorative,
                    is_chart=is_chart,
                    detailed_description=detailed_desc,
                )

            elif self.validate_alt_text and self.image_generator:
                # Image HAS alt text - validate if enabled
                existing_text = alt_text.strip()
                validated = False
                accurate = None
                issues = None
                reasoning = None
                suggested_improvement = None

                try:
                    image_path = self._extract_image_from_shape(shape)

                    if image_path:
                        if slide_context:
                            context = self._build_context_string_for_alt_text(
                                slide_context
                            )
                        else:
                            context = f"Image from slide {slide_number}"

                        result = run_async_from_sync(
                            self.image_generator.validate_alt_text(
                                image_path=image_path,
                                existing_alt_text=existing_text,
                                context=context,
                            )
                        )

                        if result.get("success"):
                            validated = True
                            accurate = result.get("is_accurate", True)
                            issues = result.get("issues", [])
                            reasoning = result.get("reasoning", "")
                            suggested_improvement = result.get("suggested_improvement")

                            if not accurate:
                                print(
                                    f"[PowerPointProcessor] Alt text validation failed for slide {slide_number}: {existing_text[:30]}... -> {issues}"
                                )

                        try:
                            os.unlink(image_path)
                        except Exception:
                            pass
                except Exception as e:
                    print(f"[PowerPointProcessor] Alt text validation failed: {e}")

                # Only return an issue if validation found problems
                if validated and not accurate:
                    return AltTextIssue(
                        slide_number=slide_number,
                        shape_id=shape.shape_id,
                        shape_name=shape.name,
                        image_type="picture",
                        has_alt_text=True,
                        existing_alt_text=existing_text,
                        suggested_alt_text=suggested_improvement,
                        alt_text_validated=True,
                        alt_text_accurate=False,
                        validation_issues=issues,
                        validation_reasoning=reasoning,
                    )

        except Exception as e:
            print(f"[PowerPointProcessor] Error checking alt text: {e}")
        return None

    def _check_text_contrast(
        self, slide_number: int, shape, slide_context: Dict = None
    ) -> Optional[ContrastIssue]:
        """Check text/background contrast ratio (with optional color blindness simulation)"""
        try:
            # Get text color
            text_frame = shape.text_frame
            if not text_frame.text.strip():
                return None

            # Try to get fill color (background)
            bg_color = self._get_fill_color(shape)
            if not bg_color:
                return None

            # Try to get text color from first run
            text_color = None
            for paragraph in text_frame.paragraphs:
                for run in paragraph.runs:
                    if run.font.color and run.font.color.rgb:
                        text_color = self._rgb_to_hex(run.font.color.rgb)
                        break
                if text_color:
                    break

            if not text_color:
                return None

            # Calculate contrast ratio
            contrast_ratio = self._calculate_contrast_ratio(text_color, bg_color)

            # Check WCAG compliance
            wcag_aa_pass = contrast_ratio >= self.wcag_aa_ratio
            wcag_aaa_pass = contrast_ratio >= self.wcag_aaa_ratio

            # Simulate color blindness if enabled
            cvd_issues = None
            if self.simulate_color_blindness and self.cvd_simulator:
                try:
                    cvd_analysis = self.cvd_simulator.analyze_color_accessibility(
                        text_color, bg_color
                    )
                    if not cvd_analysis.accessible_for_all:
                        # Convert to dict for JSON serialization
                        cvd_issues = [
                            {
                                "type": issue.color_blindness_type,
                                "contrast": issue.simulated_contrast,
                                "severity": issue.severity,
                                "description": issue.description,
                                "suggested_fix": issue.suggested_fix,
                            }
                            for issue in cvd_analysis.issues
                        ]
                except Exception as e:
                    print(
                        f"[PowerPointProcessor] Color blindness simulation failed: {e}"
                    )

            # Return issue if fails WCAG AA or has color blindness issues
            if not wcag_aa_pass or cvd_issues:
                return ContrastIssue(
                    slide_number=slide_number,
                    shape_id=shape.shape_id,
                    shape_name=shape.name,
                    foreground=text_color,
                    background=bg_color,
                    contrast_ratio=round(contrast_ratio, 2),
                    wcag_aa_pass=wcag_aa_pass,
                    wcag_aaa_pass=wcag_aaa_pass,
                    suggested_fix=self._suggest_contrast_fix(
                        text_color, bg_color, contrast_ratio
                    ),
                    color_blindness_issues=cvd_issues,
                )
        except Exception:
            pass
        return None

    def _get_fill_color(self, shape) -> Optional[str]:
        """Extract fill color from shape"""
        try:
            if shape.fill.type == 1:  # SOLID fill
                if shape.fill.fore_color.rgb:
                    return self._rgb_to_hex(shape.fill.fore_color.rgb)
        except Exception:
            pass
        return None

    def _rgb_to_hex(self, rgb) -> str:
        """Convert RGB tuple to hex string"""
        return f"#{rgb[0]:02x}{rgb[1]:02x}{rgb[2]:02x}"

    def _hex_to_rgb(self, hex_color: str) -> Tuple[int, int, int]:
        """Convert hex color to RGB tuple"""
        hex_color = hex_color.lstrip("#")
        return tuple(int(hex_color[i : i + 2], 16) for i in (0, 2, 4))

    def _calculate_contrast_ratio(self, color1: str, color2: str) -> float:
        """
        Calculate contrast ratio between two colors (WCAG formula)
        https://www.w3.org/TR/WCAG21/#dfn-contrast-ratio
        """
        rgb1 = self._hex_to_rgb(color1)
        rgb2 = self._hex_to_rgb(color2)

        l1 = self._relative_luminance(rgb1)
        l2 = self._relative_luminance(rgb2)

        lighter = max(l1, l2)
        darker = min(l1, l2)

        return (lighter + 0.05) / (darker + 0.05)

    def _relative_luminance(self, rgb: Tuple[int, int, int]) -> float:
        """Calculate relative luminance (WCAG formula)"""

        def adjust(color_value):
            color_value = color_value / 255.0
            if color_value <= 0.03928:
                return color_value / 12.92
            return ((color_value + 0.055) / 1.055) ** 2.4

        r, g, b = rgb
        r = adjust(r)
        g = adjust(g)
        b = adjust(b)

        return 0.2126 * r + 0.7152 * g + 0.0722 * b

    def _suggest_contrast_fix(
        self, text_color: str, bg_color: str, current_ratio: float
    ) -> str:
        """Suggest a fix for contrast issues"""
        needed_ratio = self.wcag_aa_ratio
        if current_ratio < needed_ratio:
            # Simple suggestion: darken text or lighten background
            text_rgb = self._hex_to_rgb(text_color)
            avg_text = sum(text_rgb) / 3

            if avg_text < 128:
                return f"Darken text color (current ratio: {current_ratio:.2f}:1, need: {needed_ratio}:1)"
            else:
                return f"Lighten background or darken text (current ratio: {current_ratio:.2f}:1, need: {needed_ratio}:1)"
        return "Contrast is sufficient"

    def _calculate_summary(
        self, slides_issues: List[SlideAccessibilityIssues]
    ) -> Dict[str, int]:
        """Calculate summary statistics"""
        contrast_count = sum(len(slide.contrast_issues) for slide in slides_issues)
        alt_text_count = sum(len(slide.alt_text_issues) for slide in slides_issues)
        title_count = sum(len(slide.title_issues) for slide in slides_issues)
        image_of_text_count = sum(
            len(slide.image_of_text_issues) for slide in slides_issues
        )

        return {
            "contrast_issues": contrast_count,
            "alt_text_issues": alt_text_count,
            "title_issues": title_count,
            "image_of_text_issues": image_of_text_count,
            "total_issues": contrast_count
            + alt_text_count
            + title_count
            + image_of_text_count,
            "slides_with_issues": sum(
                1 for slide in slides_issues if slide.total_issues > 0
            ),
        }

    def _calculate_compliance_score(
        self, summary: Dict, total_shapes: int, total_images: int
    ) -> float:
        """
        Calculate compliance score using unified scoring system.

        Uses ratio-based scoring since PPTX issues are typed (contrast, alt-text)
        rather than severity-categorized.

        - Alt text issues are treated as HIGH severity (WCAG 1.1.1 is critical)
        - Contrast issues are treated as MEDIUM severity (WCAG 1.4.3)
        - Images of text are treated as MEDIUM severity (WCAG 1.4.5)
        - Title issues are treated as MEDIUM severity (WCAG 1.3.1)
        """
        from .compliance_scoring import score_from_severity_counts

        total_elements = total_shapes + total_images
        # No early return for total_elements == 0: issues must always count
        # (issue #90 — near-empty decks scored 100.0 while listing defects).
        # Map PPTX issue types to severity
        # Missing alt text = High (blocks screen reader users)
        # Contrast issues = Medium (degraded experience)
        # Images of text = Medium (WCAG 1.4.5 AA)
        # Title issues = Medium (navigation hindrance)
        result = score_from_severity_counts(
            critical=0,
            high=summary.get("alt_text_issues", 0),
            medium=(
                summary.get("contrast_issues", 0)
                + summary.get("image_of_text_issues", 0)
                + summary.get("title_issues", 0)
            ),
            low=0,
            total_elements=total_elements,
        )
        return result.score

    def _generate_remediation_suggestions(self, summary: Dict) -> List[str]:
        """Generate high-level remediation suggestions"""
        suggestions = []

        if summary.get("title_issues", 0) > 0:
            suggestions.append(
                f"Add unique, descriptive titles to {summary['title_issues']} slides. "
                "Screen reader users rely on slide titles for navigation (WCAG 1.3.1)."
            )

        if summary["contrast_issues"] > 0:
            suggestions.append(
                f"Fix {summary['contrast_issues']} contrast ratio violations. "
                "Use darker text or lighter backgrounds to meet WCAG 2.1 AA (4.5:1 ratio)."
            )

        if summary["alt_text_issues"] > 0:
            suggestions.append(
                f"Add descriptive alt text to {summary['alt_text_issues']} images. "
                "Describe the content and purpose of each image for screen reader users."
            )

        if summary.get("image_of_text_issues", 0) > 0:
            suggestions.append(
                f"Convert {summary['image_of_text_issues']} images of text to real text. "
                "Text in images cannot be resized, searched, or read by screen readers (WCAG 1.4.5)."
            )

        if summary["total_issues"] == 0:
            suggestions.append(
                "✅ No accessibility issues detected. Presentation is WCAG 2.1 compliant!"
            )

        return suggestions

    def process_directory(self, directory: str) -> List[PowerPointProcessingResult]:
        """
        Batch process all PPTX files in a directory

        Args:
            directory: Path to directory containing PPTX files

        Returns:
            List of PowerPointProcessingResult for each file
        """
        results = []
        pptx_files = list(Path(directory).glob("*.pptx"))

        for pptx_file in pptx_files:
            try:
                result = self.process_pptx(str(pptx_file))
                results.append(result)
            except Exception as e:
                print(f"Error processing {pptx_file}: {e}")

        return results
