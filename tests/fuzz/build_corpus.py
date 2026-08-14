"""Build the generator-diverse fuzzing corpus (issue #88).

Both 2026-08-12 prod bugs were code assuming well-formed happy-path output
from a generator we don't control (Gemini thinking-token truncation; DOCX with
no default paragraph style). Real university uploads are dominated by exactly
those generators, so the corpus varies the GENERATOR, not the content.

Generators: pandoc (docx/pptx/pdf via pdflatex), python-docx, python-pptx,
openpyxl, hand-stripped OOXML variants. When LibreOffice (soffice) is on PATH
it round-trips the pandoc DOCX through LibreOffice as an extra generator.
Google Docs exports cannot be scripted server-side — drop manual exports into
corpus/ before running for that column.

Usage:  python tests/fuzz/build_corpus.py   (from backend/, venv active)
"""

import re
import shutil
import subprocess
import zipfile
from pathlib import Path

HERE = Path(__file__).parent
CORPUS = HERE / "corpus"

SOURCE_MD = """% AZ-900 Study Notes
% Course Staff

# Domain 1 — Cloud Concepts

Cloud computing delivers services over the internet.

## Benefits

- High availability
- Scalability
- Elasticity

## Cost model

| Model | Commitment | Discount |
|-------|-----------|----------|
| Pay-as-you-go | None | 0% |
| Reserved | 1-3 years | up to 72% |

![Treatment outcomes chart](test-chart.png)

# Domain 2 — Architecture

Regions contain availability zones.

```python
def deploy(region):
    return f"deployed to {region}"
```

1. Pick a region
2. Pick a zone
3. Deploy
"""


def make_chart(path: Path) -> None:
    """Simple bar-chart-like PNG via Pillow (already a backend dependency)."""
    from PIL import Image, ImageDraw

    img = Image.new("RGB", (640, 400), "white")
    d = ImageDraw.Draw(img)
    d.text((180, 10), "Treatment outcomes by group", fill="black")
    for i, (h, c) in enumerate(
        [(120, "gray"), (170, "steelblue"), (260, "seagreen"), (190, "indianred")]
    ):
        x = 80 + i * 130
        d.rectangle([x, 360 - h, x + 90, 360], fill=c)
    d.line([(60, 360), (600, 360)], fill="black", width=2)
    img.save(path)


def strip_default_styles(src: Path, dst: Path) -> None:
    """Remove every w:default flag from word/styles.xml.

    Reproduces the class of file behind the 2026-08-12 prod crash: valid OOXML
    that merely omits something Word always writes.
    """
    with zipfile.ZipFile(src) as zin, zipfile.ZipFile(dst, "w") as zout:
        for item in zin.infolist():
            data = zin.read(item.filename)
            if item.filename == "word/styles.xml":
                data = re.sub(rb'\sw:default="(?:1|true)"', b"", data)
            zout.writestr(item, data)


def main() -> None:
    CORPUS.mkdir(parents=True, exist_ok=True)
    chart = HERE / "test-chart.png"
    make_chart(chart)
    (HERE / "source.md").write_text(SOURCE_MD)

    if shutil.which("pandoc"):

        def pandoc(*args):
            subprocess.run(["pandoc", *args], cwd=HERE, check=True, timeout=180)

        pandoc("source.md", "-o", str(CORPUS / "pandoc.docx"))
        pandoc("source.md", "-o", str(CORPUS / "pandoc.pptx"))
        try:
            pandoc("source.md", "-o", str(CORPUS / "pandoc.pdf"))  # needs pdflatex
        except subprocess.CalledProcessError:
            print("pandoc.pdf skipped (no LaTeX engine)")
        strip_default_styles(CORPUS / "pandoc.docx", CORPUS / "pandoc-stripped.docx")
    else:
        print("pandoc not on PATH — pandoc generators skipped")

    import docx

    d = docx.Document()
    d.add_heading("Domain 1 — Cloud Concepts", level=1)
    d.add_paragraph("Cloud computing delivers services over the internet.")
    d.add_heading("Benefits", level=2)
    for item in ["High availability", "Scalability", "Elasticity"]:
        d.add_paragraph(item, style="List Bullet")
    d.add_picture(str(chart))  # no alt text — should be flagged, not crash
    t = d.add_table(rows=2, cols=2)
    t.rows[0].cells[0].text = "Model"
    t.rows[0].cells[1].text = "Discount"
    t.rows[1].cells[0].text = "Reserved"
    t.rows[1].cells[1].text = "72%"
    d.save(str(CORPUS / "pythondocx.docx"))
    strip_default_styles(
        CORPUS / "pythondocx.docx", CORPUS / "pythondocx-stripped.docx"
    )

    empty = docx.Document()
    empty.save(str(CORPUS / "pythondocx-empty.docx"))

    # Large document — the only corpus file above MIN_ELEMENTS_FOR_RATIO, so
    # it is what keeps the ratio-scoring branch covered (found 2026-08-12:
    # every other file lands on the penalty path, leaving ratio scoring with
    # zero fuzz signal).
    large = docx.Document()
    large.add_heading("Full Semester Notes", level=1)
    for week in range(1, 11):
        large.add_heading(f"Week {week}", level=2)
        for para in range(4):
            large.add_paragraph(
                f"Week {week} content paragraph {para + 1}: cloud concepts, "
                "architecture, pricing, and governance in depth."
            )
        large.add_paragraph(f"Week {week} reading", style="List Bullet")
    large.add_picture(str(chart))  # no alt text — one issue over many elements
    large.save(str(CORPUS / "pythondocx-large.docx"))

    import pptx

    p = pptx.Presentation()
    slide = p.slides.add_slide(p.slide_layouts[6])  # blank layout: no title
    box = slide.shapes.add_textbox(0, 0, 100000, 100000)
    box.text_frame.text = "Untitled slide with a lone textbox"
    slide.shapes.add_picture(str(chart), 0, 200000)  # no alt text
    p.save(str(CORPUS / "pythonpptx.pptx"))

    import openpyxl

    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = "Grades"
    ws.append(["Student", "Score"])
    ws.append(["A", 91])
    ws.merge_cells("A4:B5")
    ws["A4"] = "merged block"
    wb.create_sheet("Empty Sheet")
    wb.save(str(CORPUS / "openpyxl.xlsx"))

    soffice = shutil.which("soffice") or (
        "/Applications/LibreOffice.app/Contents/MacOS/soffice"
        if Path("/Applications/LibreOffice.app").exists()
        else None
    )
    if soffice and (CORPUS / "pandoc.docx").exists():
        for fmt in ("docx", "pdf"):
            subprocess.run(
                [
                    soffice,
                    "--headless",
                    "--convert-to",
                    fmt,
                    str(CORPUS / "pandoc.docx"),
                    "--outdir",
                    str(CORPUS / "libreoffice"),
                ],
                check=True,
                timeout=300,
            )
        for f in (CORPUS / "libreoffice").iterdir():
            f.rename(CORPUS / f"libreoffice.{f.suffix.lstrip('.')}")
        (CORPUS / "libreoffice").rmdir()
    else:
        print("LibreOffice not found — libreoffice generators skipped")

    print("corpus files:")
    for f in sorted(CORPUS.iterdir()):
        print(f"  {f.name}  ({f.stat().st_size} bytes)")


if __name__ == "__main__":
    main()
