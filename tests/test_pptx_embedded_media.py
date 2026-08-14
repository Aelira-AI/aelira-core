"""
Tests for Embedded Media Checks (Task 9)

Tests cover:
- EmbeddedMedia model
- EmbeddedMediaIssue model
- Video/audio detection in PPTX files
- Caption/subtitle checking
- Transcript checking in speaker notes
"""

import pytest
import tempfile
import os
import zipfile

from src.education.pptx_processor import (
    PowerPointProcessor,
    EmbeddedMedia,
    EmbeddedMediaIssue,
)


class TestEmbeddedMediaModel:
    """Test EmbeddedMedia Pydantic model."""

    def test_embedded_media_creation(self):
        """Test creating an embedded media object."""
        media = EmbeddedMedia(
            slide_number=1,
            media_index=0,
            media_type="video",
            file_name="presentation_video.mp4",
            file_path="ppt/media/media1.mp4",
            duration_seconds=120.5,
            has_captions=True,
            has_transcript=True,
            content_type="video/mp4",
        )

        assert media.slide_number == 1
        assert media.media_type == "video"
        assert media.has_captions is True
        assert media.duration_seconds == 120.5

    def test_embedded_media_audio(self):
        """Test embedded audio media."""
        media = EmbeddedMedia(
            slide_number=2,
            media_index=0,
            media_type="audio",
            file_name="narration.mp3",
            has_captions=False,
            has_transcript=False,
        )

        assert media.media_type == "audio"
        assert media.has_captions is False

    def test_embedded_media_minimal(self):
        """Test embedded media with minimal fields."""
        media = EmbeddedMedia(
            slide_number=1,
            media_index=0,
            media_type="video",
        )

        assert media.file_name is None
        assert media.duration_seconds is None
        assert media.has_captions is False


class TestEmbeddedMediaIssueModel:
    """Test EmbeddedMediaIssue Pydantic model."""

    def test_media_issue_creation(self):
        """Test creating an embedded media issue."""
        issue = EmbeddedMediaIssue(
            slide_number=1,
            media_index=0,
            media_type="video",
            file_name="video.mp4",
            issue_type="missing_captions",
            recommendations=[
                "Add synchronized captions to the video",
                "Consider automatic caption services",
            ],
            suggested_fix="Add synchronized captions to the video (WCAG 1.2.2)",
        )

        assert issue.slide_number == 1
        assert issue.media_type == "video"
        assert issue.issue_type == "missing_captions"
        assert len(issue.recommendations) == 2

    def test_media_issue_audio(self):
        """Test embedded audio issue."""
        issue = EmbeddedMediaIssue(
            slide_number=3,
            media_index=0,
            media_type="audio",
            file_name="narration.mp3",
            issue_type="missing_transcript",
            recommendations=["Provide a text transcript for the audio content"],
            suggested_fix="Add transcript to speaker notes",
        )

        assert issue.media_type == "audio"
        assert issue.issue_type == "missing_transcript"


class TestEmbeddedMediaDetection:
    """Test embedded media detection in PPTX files."""

    @pytest.fixture
    def processor(self):
        """Create PowerPointProcessor for testing."""
        return PowerPointProcessor()

    def test_check_embedded_media_empty_slide(self, processor):
        """Test media check on slide without media."""
        from pptx import Presentation

        with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as f:
            prs = Presentation()
            slide_layout = prs.slide_layouts[5]  # Blank
            slide = prs.slides.add_slide(slide_layout)
            prs.save(f.name)

            issues = processor._check_embedded_media(slide, 1, f.name)

            # Should return empty list for slide without media
            assert isinstance(issues, list)

        os.unlink(f.name)

    def test_check_embedded_media_returns_list(self, processor):
        """Test that check_embedded_media always returns a list."""
        from pptx import Presentation

        with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as f:
            prs = Presentation()
            slide_layout = prs.slide_layouts[5]
            slide = prs.slides.add_slide(slide_layout)
            prs.save(f.name)

            result = processor._check_embedded_media(slide, 1, f.name)

            assert isinstance(result, list)

        os.unlink(f.name)


class TestSpeakerNotesTranscript:
    """Test transcript detection in speaker notes."""

    @pytest.fixture
    def processor(self):
        """Create PowerPointProcessor for testing."""
        return PowerPointProcessor()

    def test_check_notes_no_transcript(self, processor):
        """Test when speaker notes have no transcript."""
        from pptx import Presentation

        with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as f:
            prs = Presentation()
            slide_layout = prs.slide_layouts[5]
            slide = prs.slides.add_slide(slide_layout)
            prs.save(f.name)

            with zipfile.ZipFile(f.name, "r") as zf:
                result = processor._check_speaker_notes_for_transcript(slide, 1, zf)

            # Should return False when no notes
            assert result is False

        os.unlink(f.name)


class TestCaptionDetection:
    """Test caption/subtitle detection."""

    @pytest.fixture
    def processor(self):
        """Create PowerPointProcessor for testing."""
        return PowerPointProcessor()

    def test_check_captions_none(self, processor):
        """Test when no captions exist."""
        from pptx import Presentation

        with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as f:
            prs = Presentation()
            slide_layout = prs.slide_layouts[5]
            prs.slides.add_slide(slide_layout)
            prs.save(f.name)

            with zipfile.ZipFile(f.name, "r") as zf:
                result = processor._check_for_captions(zf, "ppt/media/video.mp4")

            assert result is False

        os.unlink(f.name)

    def test_check_captions_with_vtt(self, processor):
        """Test detection of WebVTT caption file."""
        from pptx import Presentation

        with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as f:
            prs = Presentation()
            slide_layout = prs.slide_layouts[5]
            prs.slides.add_slide(slide_layout)
            prs.save(f.name)

            # Add fake caption file
            with zipfile.ZipFile(f.name, "a") as zf:
                zf.writestr(
                    "ppt/media/video.vtt", "WEBVTT\n\n00:00.000 --> 00:01.000\nHello"
                )

            with zipfile.ZipFile(f.name, "r") as zf:
                result = processor._check_for_captions(zf, "ppt/media/video.mp4")

            # Should detect the VTT file
            assert result is True

        os.unlink(f.name)

    def test_check_captions_with_srt(self, processor):
        """Test detection of SRT caption file."""
        from pptx import Presentation

        with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as f:
            prs = Presentation()
            slide_layout = prs.slide_layouts[5]
            prs.slides.add_slide(slide_layout)
            prs.save(f.name)

            # Add fake SRT caption file
            with zipfile.ZipFile(f.name, "a") as zf:
                zf.writestr(
                    "ppt/media/presentation.srt",
                    "1\n00:00:00,000 --> 00:00:01,000\nHello",
                )

            with zipfile.ZipFile(f.name, "r") as zf:
                result = processor._check_for_captions(zf, "ppt/media/presentation.mp4")

            # Should detect the SRT file
            assert result is True

        os.unlink(f.name)


class TestPptxProcessorIntegration:
    """Test embedded media integration with PowerPointProcessor."""

    @pytest.fixture
    def processor(self):
        """Create PowerPointProcessor for testing."""
        return PowerPointProcessor()

    def test_process_pptx_includes_media_issues(self, processor):
        """Test that process_pptx includes embedded_media_issues in slides."""
        from pptx import Presentation

        with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as f:
            prs = Presentation()
            slide_layout = prs.slide_layouts[5]
            prs.slides.add_slide(slide_layout)
            prs.save(f.name)

            result = processor.process_pptx(f.name)

            # Each slide should have embedded_media_issues field
            assert len(result.slides) > 0
            for slide in result.slides:
                assert hasattr(slide, "embedded_media_issues")
                assert isinstance(slide.embedded_media_issues, list)

        os.unlink(f.name)

    def test_media_issues_counted_in_total(self, processor):
        """Test that media issues are counted in total issues."""
        from pptx import Presentation

        with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as f:
            prs = Presentation()
            slide_layout = prs.slide_layouts[5]
            prs.slides.add_slide(slide_layout)
            prs.save(f.name)

            result = processor.process_pptx(f.name)

            for slide in result.slides:
                media_count = len(slide.embedded_media_issues)
                # Total should include media issues
                if media_count > 0:
                    assert slide.total_issues >= media_count

        os.unlink(f.name)


class TestSlideAccessibilityIssuesModel:
    """Test SlideAccessibilityIssues includes embedded_media_issues."""

    def test_slide_issues_has_media_field(self):
        """Test that SlideAccessibilityIssues has embedded_media_issues field."""
        from src.education.pptx_processor import SlideAccessibilityIssues

        slide_issues = SlideAccessibilityIssues(
            slide_number=1,
            slide_title="Test Slide",
            contrast_issues=[],
            alt_text_issues=[],
            title_issues=[],
            image_of_text_issues=[],
            animation_issues=[],
            embedded_media_issues=[],
            total_issues=0,
        )

        assert hasattr(slide_issues, "embedded_media_issues")
        assert slide_issues.embedded_media_issues == []

    def test_slide_issues_with_media_issues(self):
        """Test SlideAccessibilityIssues with media issues."""
        from src.education.pptx_processor import SlideAccessibilityIssues

        media_issue = EmbeddedMediaIssue(
            slide_number=1,
            media_index=0,
            media_type="video",
            file_name="video.mp4",
            issue_type="missing_captions",
            recommendations=["Add captions"],
            suggested_fix="Add captions",
        )

        slide_issues = SlideAccessibilityIssues(
            slide_number=1,
            slide_title="Test",
            contrast_issues=[],
            alt_text_issues=[],
            embedded_media_issues=[media_issue],
            total_issues=1,
        )

        assert len(slide_issues.embedded_media_issues) == 1
        assert slide_issues.embedded_media_issues[0].media_type == "video"


if __name__ == "__main__":
    pytest.main([__file__, "-v", "-s"])
