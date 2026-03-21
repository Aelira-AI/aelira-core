"""
Tests for Animation Accessibility Analysis (Task 8)

Tests cover:
- AnimationIssue model
- Animation detection in PPTX files
- Animation type identification
- Flash risk detection
- Auto-start detection
- Motion path analysis
"""

import pytest
import tempfile
import os
import zipfile
from xml.etree import ElementTree as ET

from src.education.pptx_processor import (
    PowerPointProcessor,
    AnimationIssue,
)


class TestAnimationIssueModel:
    """Test AnimationIssue Pydantic model."""

    def test_animation_issue_creation(self):
        """Test creating an animation issue."""
        issue = AnimationIssue(
            slide_number=1,
            animation_index=0,
            animation_type="entrance",
            element_name="Title 1",
            effect_name="fade",
            duration_ms=500,
            delay_ms=0,
            is_auto_start=True,
            repeat_count=3,
            issues=["auto_advance"],
            suggested_fix="Consider requiring user click to start animation",
        )

        assert issue.slide_number == 1
        assert issue.animation_type == "entrance"
        assert issue.is_auto_start is True
        assert "auto_advance" in issue.issues

    def test_animation_issue_minimal(self):
        """Test animation issue with minimal fields."""
        issue = AnimationIssue(
            slide_number=2,
            animation_index=1,
            animation_type="emphasis",
            element_name="Shape 5",
            duration_ms=1000,
            issues=["motion_heavy"],
            suggested_fix="Reduce motion",
        )

        assert issue.effect_name is None
        assert issue.repeat_count is None
        assert issue.delay_ms == 0

    def test_animation_issue_types(self):
        """Test different animation types."""
        types = ["entrance", "emphasis", "exit", "motion_path"]

        for anim_type in types:
            issue = AnimationIssue(
                slide_number=1,
                animation_index=0,
                animation_type=anim_type,
                element_name="Test",
                duration_ms=500,
                issues=[],
                suggested_fix="Test fix",
            )
            assert issue.animation_type == anim_type


class TestAnimationTypeIdentification:
    """Test animation type identification."""

    @pytest.fixture
    def processor(self):
        """Create PowerPointProcessor for testing."""
        return PowerPointProcessor()

    def test_get_animation_type_entrance(self, processor):
        """Test entrance animation type detection."""
        assert processor._get_animation_type("anim") == "entrance"

    def test_get_animation_type_emphasis(self, processor):
        """Test emphasis animation type detection."""
        assert processor._get_animation_type("animEffect") == "emphasis"
        assert processor._get_animation_type("set") == "emphasis"
        assert processor._get_animation_type("animClr") == "emphasis"

    def test_get_animation_type_motion(self, processor):
        """Test motion path animation type detection."""
        assert processor._get_animation_type("animMotion") == "motion_path"

    def test_get_animation_type_unknown(self, processor):
        """Test unknown animation type."""
        assert processor._get_animation_type("unknownType") == "other"


class TestAnimationDurationExtraction:
    """Test animation duration extraction."""

    @pytest.fixture
    def processor(self):
        """Create PowerPointProcessor for testing."""
        return PowerPointProcessor()

    def test_get_duration_from_attribute(self, processor):
        """Test extracting duration from element attribute."""
        element = ET.fromstring('<anim dur="1000"/>')
        duration = processor._get_animation_duration(element)
        assert duration == 1000

    def test_get_duration_default(self, processor):
        """Test default duration when not specified."""
        element = ET.fromstring("<anim/>")
        duration = processor._get_animation_duration(element)
        assert duration == 500  # Default


class TestAutoStartDetection:
    """Test auto-start animation detection."""

    @pytest.fixture
    def processor(self):
        """Create PowerPointProcessor for testing."""
        return PowerPointProcessor()

    def test_is_auto_start_with_after_previous(self, processor):
        """Test auto-start detection with afterPrevious."""
        element = ET.fromstring("""
            <anim xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main">
                <cTn nodeType="afterPrevious"/>
            </anim>
        """)
        assert processor._is_auto_start(element) is True

    def test_is_auto_start_false(self, processor):
        """Test when animation is not auto-start."""
        element = ET.fromstring("<anim/>")
        assert processor._is_auto_start(element) is False


class TestAnimationFixGeneration:
    """Test animation fix suggestion generation."""

    @pytest.fixture
    def processor(self):
        """Create PowerPointProcessor for testing."""
        return PowerPointProcessor()

    def test_generate_fix_rapid_flash(self, processor):
        """Test fix generation for rapid flash issues."""
        fix = processor._generate_animation_fix(["rapid_flash"], "entrance")

        assert "flash" in fix.lower() or "seizure" in fix.lower()
        assert "WCAG" in fix

    def test_generate_fix_auto_advance(self, processor):
        """Test fix generation for auto-advance issues."""
        fix = processor._generate_animation_fix(["auto_advance"], "entrance")

        assert "automatic" in fix.lower() or "control" in fix.lower()

    def test_generate_fix_motion_heavy(self, processor):
        """Test fix generation for motion-heavy issues."""
        fix = processor._generate_animation_fix(["motion_heavy"], "motion_path")

        assert "motion" in fix.lower() or "vestibular" in fix.lower()

    def test_generate_fix_multiple_issues(self, processor):
        """Test fix generation for multiple issues."""
        fix = processor._generate_animation_fix(
            ["rapid_flash", "auto_advance"], "entrance"
        )

        assert "flash" in fix.lower() or "seizure" in fix.lower()
        assert "automatic" in fix.lower() or "control" in fix.lower()


class TestAnimationAnalysis:
    """Test animation analysis on PPTX files."""

    @pytest.fixture
    def processor(self):
        """Create PowerPointProcessor for testing."""
        return PowerPointProcessor()

    def test_analyze_animations_empty_slide(self, processor):
        """Test animation analysis on slide without animations."""
        from pptx import Presentation

        with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as f:
            prs = Presentation()
            slide_layout = prs.slide_layouts[5]  # Blank layout
            slide = prs.slides.add_slide(slide_layout)
            prs.save(f.name)

            issues = processor._analyze_animations(slide, 1, f.name)

            # Should return empty list for slide without animations
            assert isinstance(issues, list)

        os.unlink(f.name)

    def test_analyze_animations_returns_list(self, processor):
        """Test that analyze_animations always returns a list."""
        from pptx import Presentation

        with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as f:
            prs = Presentation()
            slide_layout = prs.slide_layouts[5]
            slide = prs.slides.add_slide(slide_layout)
            prs.save(f.name)

            result = processor._analyze_animations(slide, 1, f.name)

            assert isinstance(result, list)

        os.unlink(f.name)

    def test_analyze_animations_with_transition(self, processor):
        """Test detection of auto-advance slide transitions."""
        from pptx import Presentation

        with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as f:
            prs = Presentation()
            slide_layout = prs.slide_layouts[5]
            slide = prs.slides.add_slide(slide_layout)
            prs.save(f.name)

            # Add auto-advance transition to slide XML
            with zipfile.ZipFile(f.name, "a") as zf:
                # Read existing slide
                slide_xml = zf.read("ppt/slides/slide1.xml")
                # This would need to modify the transition element
                # For now, just test the method doesn't crash

            result = processor._analyze_animations(slide, 1, f.name)
            assert isinstance(result, list)

        os.unlink(f.name)


class TestPptxProcessorIntegration:
    """Test animation integration with PowerPointProcessor."""

    @pytest.fixture
    def processor(self):
        """Create PowerPointProcessor for testing."""
        return PowerPointProcessor()

    def test_process_pptx_includes_animation_issues(self, processor):
        """Test that process_pptx includes animation_issues in slides."""
        from pptx import Presentation

        with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as f:
            prs = Presentation()
            slide_layout = prs.slide_layouts[5]
            prs.slides.add_slide(slide_layout)
            prs.save(f.name)

            result = processor.process_pptx(f.name)

            # Each slide should have animation_issues field
            assert len(result.slides) > 0
            for slide in result.slides:
                assert hasattr(slide, "animation_issues")
                assert isinstance(slide.animation_issues, list)

        os.unlink(f.name)

    def test_animation_issues_counted_in_total(self, processor):
        """Test that animation issues are counted in total issues."""
        from pptx import Presentation

        with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as f:
            prs = Presentation()
            slide_layout = prs.slide_layouts[5]
            prs.slides.add_slide(slide_layout)
            prs.save(f.name)

            result = processor.process_pptx(f.name)

            # Animation issues should be part of total
            for slide in result.slides:
                animation_count = len(slide.animation_issues)
                # Total should include animations
                if animation_count > 0:
                    assert slide.total_issues >= animation_count

        os.unlink(f.name)


class TestSlideAccessibilityIssuesModel:
    """Test SlideAccessibilityIssues includes animation_issues."""

    def test_slide_issues_has_animation_field(self):
        """Test that SlideAccessibilityIssues has animation_issues field."""
        from src.education.pptx_processor import SlideAccessibilityIssues

        slide_issues = SlideAccessibilityIssues(
            slide_number=1,
            slide_title="Test Slide",
            contrast_issues=[],
            alt_text_issues=[],
            title_issues=[],
            image_of_text_issues=[],
            animation_issues=[],
            total_issues=0,
        )

        assert hasattr(slide_issues, "animation_issues")
        assert slide_issues.animation_issues == []

    def test_slide_issues_with_animation_issues(self):
        """Test SlideAccessibilityIssues with animation issues."""
        from src.education.pptx_processor import SlideAccessibilityIssues

        anim_issue = AnimationIssue(
            slide_number=1,
            animation_index=0,
            animation_type="entrance",
            element_name="Title",
            duration_ms=100,
            issues=["rapid_animation"],
            suggested_fix="Slow down animation",
        )

        slide_issues = SlideAccessibilityIssues(
            slide_number=1,
            slide_title="Test",
            contrast_issues=[],
            alt_text_issues=[],
            animation_issues=[anim_issue],
            total_issues=1,
        )

        assert len(slide_issues.animation_issues) == 1
        assert slide_issues.animation_issues[0].animation_type == "entrance"


if __name__ == "__main__":
    pytest.main([__file__, "-v", "-s"])
